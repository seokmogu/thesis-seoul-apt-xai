#!/usr/bin/env python3
"""v7 발표 PPT 완전 재작성 (깨끗한 버전):
- 표지·목차·섹션표지·참고문헌·감사 슬라이드는 원본 디자인 유지
- content 슬라이드는 모든 shape 제거 후 새로 작성
- Pretendard 폰트 강제
- 내용 overflow 방지 위해 영역 크기 조절
"""
import os, shutil
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = '/Users/seokmogu/Downloads/양재영_중간발표_1203.pptx'
DST = '/Users/seokmogu/project/thesis-seoul-apt-xai/paper/중간발표_v7_m2price.pptx'
PLOTS = '/Users/seokmogu/project/thesis-seoul-apt-xai/results/plots_v7_m2price'

shutil.copy(SRC, DST)
prs = Presentation(DST)

CYAN = RGBColor(0x08, 0xA5, 0xC1)
DARK = RGBColor(0x39, 0x39, 0x39)
LIGHT = RGBColor(0x7A, 0x7A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_KR = 'Pretendard'

W = prs.slide_width
H = prs.slide_height

# ========== 유틸 ==========
def remove_shape(shape):
    try: shape._element.getparent().remove(shape._element)
    except: pass

def strip_all(slide, keep_types=()):
    """keep_types 외 모든 shape 제거."""
    for sh in list(slide.shapes):
        if type(sh).__name__ in keep_types:
            continue
        remove_shape(sh)

def add_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_font_ea(run, fontname=FONT_KR):
    rPr = run._r.get_or_add_rPr()
    for ex in rPr.findall(qn('a:ea')): rPr.remove(ex)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', fontname)
    for ex in rPr.findall(qn('a:latin')): rPr.remove(ex)
    la = etree.SubElement(rPr, qn('a:latin'))
    la.set('typeface', fontname)

def add_p(tf, text, size=14, bold=False, color=DARK, align=None, first=False):
    if first:
        p = tf.paragraphs[0]
        # 기존 내용 지우기
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    else:
        p = tf.add_paragraph()
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT_KR
    r.font.size = Pt(size)
    r.font.bold = bold
    try: r.font.color.rgb = color
    except: pass
    set_font_ea(r)
    return r

def write_content_slide(slide, page_num, section_num, section_label, title, lines,
                        font_size=15, title_size=30):
    """content slide 완전 재작성: 모든 shape 제거 → 절 라벨 + 제목 + 본문 + 페이지번호 재구성."""
    strip_all(slide)
    # 절 라벨 (좌상)
    label_box = add_box(slide, Inches(0.8), Inches(0.7), Inches(3.5), Inches(0.6))
    tf = label_box.text_frame
    tf.word_wrap = False
    add_p(tf, f'{section_num}  {section_label}', size=18, bold=False, color=LIGHT, first=True)
    # 큰 제목
    title_box = add_box(slide, Inches(0.8), Inches(1.2), Inches(18), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    add_p(tf, title, size=title_size, bold=True, color=DARK, first=True)
    # 본문
    body_box = add_box(slide, Inches(0.8), Inches(2.3), W - Inches(1.6), H - Inches(3.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        sz = font_size
        bold = False
        color = DARK
        if ln == '':
            add_p(tf, ' ', size=8, first=(i==0))
            continue
        if ln.startswith('[') and ln.endswith(']'):
            sz = font_size + 3; bold = True; color = CYAN
        elif ln.startswith('→'):
            sz = font_size; bold = True; color = CYAN
        elif ln.startswith('※') or ln.startswith('*'):
            sz = font_size - 2; color = LIGHT
        add_p(tf, ln, size=sz, bold=bold, color=color, first=(i==0))
    # 페이지 번호 (우하)
    pn = add_box(slide, W - Inches(1.5), H - Inches(0.8), Inches(1.2), Inches(0.5))
    tf = pn.text_frame
    add_p(tf, str(page_num), size=14, color=LIGHT, align=PP_ALIGN.RIGHT, first=True)

def fix_font_eastasian(slide):
    """기존 슬라이드의 한글 run에 Pretendard 강제."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        has_kr = any('가' <= c <= '힣' for c in r.text)
                        if has_kr:
                            r.font.name = FONT_KR
                            set_font_ea(r)

# ========== 슬라이드 1 (표지) ==========
s = prs.slides[0]
strip_all(s)  # 원본 양재영 표지 전부 제거하고 새로 만들기
# 어두운 배경은 슬라이드 마스터 유지 — 만약 없으면 추가
# 제목
title_box = add_box(s, Inches(1.5), Inches(3.3), Inches(17), Inches(2.5))
tf = title_box.text_frame; tf.word_wrap = True
add_p(tf, 'XGBoost와 SHAP을 활용한', size=36, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER, first=True)
add_p(tf, '서울시 아파트 단위면적당 매매가격의', size=40, bold=True,
      color=CYAN, align=PP_ALIGN.CENTER)
add_p(tf, '설명 패턴 분석', size=40, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
# 부제
sub_box = add_box(s, Inches(1.5), Inches(6.5), Inches(17), Inches(0.8))
tf = sub_box.text_frame
add_p(tf, '단위면적 정규화 · 행정동 분석 · 지역 별도 모형', size=22,
      color=WHITE, align=PP_ALIGN.CENTER, first=True)
# 정보
info_box = add_box(s, Inches(2), Inches(8.5), Inches(16), Inches(2))
tf = info_box.text_frame
add_p(tf, '한양대학교 부동산융합대학원 · 도시부동산정책전공', size=18,
      color=WHITE, align=PP_ALIGN.CENTER, first=True)
add_p(tf, '석사학위 중간발표  |  박 현 근', size=24, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER)
add_p(tf, '2026. 04.', size=16, color=WHITE, align=PP_ALIGN.CENTER)

# ========== 슬라이드 2 (목차) ==========
# 원본 목차 디자인 유지, 라벨만 교체
s = prs.slides[1]
toc_map = {
    '서론': '서 론',
    '이론': '이론적 배경 및 선행연구',
    '분석의': '연구 설계 및 방법',
    '분석결과': '실증 분석 결과',
    '결론': '결론 및 시사점',
}
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() in toc_map:
                    r.text = toc_map[r.text.strip()]
                    r.font.name = FONT_KR
                    set_font_ea(r)
fix_font_eastasian(s)

# ========== 섹션 표지 슬라이드 ==========
# Slide 3 (01 서론), 8 (02), 18 (03), 26 (04), 32 (05)
SECTION_LABELS = {
    2: ('01', '서 론'),
    7: ('02', '이론적 배경 및 선행연구'),
    17: ('03', '연구 설계 및 방법'),
    25: ('04', '실증 분석 결과'),
    31: ('05', '결론 및 시사점'),
}
SECTION_SUBITEMS = {
    2: ['1. 연구의 배경 및 목적', '2. 연구의 범위 및 방법'],
    7: ['1. 헤도닉 가격모형과 정규화', '2. 머신러닝 기반 예측',
        '3. XAI와 SHAP', '4. 선행연구 종합'],
    17: ['1. 연구 흐름', '2. 변수 설정 및 분할', '3. 모형·평가지표',
         '4. SHAP 분석', '5. 강건성·지역 별도 모형'],
    25: ['1. 기술통계 및 지역별 분포', '2. 상관관계 및 VIF',
         '3. Ablation·Moran\'s I', '4. 모형 성능', '5. SHAP·지역 비교'],
    31: ['1. 시사점', '2. 한계 및 향후 과제'],
}
for idx, (num, label) in SECTION_LABELS.items():
    s = prs.slides[idx]
    # 원본 모든 텍스트 박스 제거 (배경 곡선 그래픽 + 숫자는 유지 시도)
    # 단, shape_type이 TEXT_BOX이고 text가 있는 것만 제거
    from pptx.util import Inches as I
    to_remove = []
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            to_remove.append(sh)
    for sh in to_remove:
        remove_shape(sh)
    # 신규 작성: 좌측에 큰 숫자, 우측에 레이블 + 하위 항목
    # 큰 섹션 번호 (왼쪽)
    num_box = s.shapes.add_textbox(I(0.8), I(3.5), I(5), I(3.5))
    tf = num_box.text_frame
    add_p(tf, num, size=250, bold=False, color=CYAN, first=True)
    # 섹션 제목 (오른쪽)
    title_box = s.shapes.add_textbox(I(6.5), I(4.0), I(12), I(1.8))
    tf = title_box.text_frame
    add_p(tf, label, size=58, bold=True, color=WHITE, first=True)
    # 하위 항목 목록
    sub = SECTION_SUBITEMS.get(idx, [])
    if sub:
        sub_box = s.shapes.add_textbox(I(6.5), I(6.2), I(12), I(4))
        tf = sub_box.text_frame
        for i, item in enumerate(sub):
            add_p(tf, item, size=22, color=WHITE, first=(i==0))
    fix_font_eastasian(s)

# ========== content 슬라이드 작성 ==========
# (slide 1-based, section num, section label, title, lines)
CONTENTS = {
    4: ('1절', '서 론', '연구의 배경', [
        '[한국 부동산 시장의 구조적 중요성]',
        '• 한국 가계 자산의 75.2%가 실물자산, 대부분 부동산 (통계청, 2024)',
        '• 서울 아파트 = 자산 축적·주거 안정·거시경제 파급효과의 핵심 매개',
        '',
        '[기존 가격 예측 접근의 한계]',
        '• 전통적 OLS 헤도닉 모형: 선형 가정 → 비선형·상호작용 포착 한계',
        '• 머신러닝(RF·XGBoost): 예측 성능 우수하나 "블랙박스" 문제',
        '',
        '[XAI의 등장과 국내 공백]',
        '• SHAP (Lundberg & Lee, 2017): 예측값 기여도 이론적 분해',
        '• 해외에서 SHAP 기반 부동산 연구 축적',
        '• 국내 아파트 시장에 XAI를 체계적으로 적용한 연구는 여전히 부족',
    ]),
    5: ('1절', '서 론', '연구의 필요성 및 목적', [
        '[연구 동기 — 왜 이 주제인가]',
        '• AVM(자동감정평가) 실무 확산 ↔ 국내 설명 프레임워크 부재',
        '• 기존 한국 ML 연구의 3공백:',
        '   ① 자치구(25개) 단위 집계   ② 단일 모형 예측 비교   ③ 통합 SHAP만 보고',
        '• 총가격 모형: 면적이 SHAP 상위 독점 → 질적 신호 가려짐',
        '',
        '[본 연구의 3가지 차별성]',
        '① 단위면적당 가격 정규화 = log(거래금액/전용면적) → 규모효과 분리',
        '② 행정동(215개) 단위 세분화 → 기존 자치구(25개) 대비 공간 정밀도',
        '③ 지역 별도 모형 3개 (전체 · 강남3구 · 비강남 22개구) 직접 비교',
        '',
        '[연구 목적]',
        '→ 서울 아파트 단위가격 결정의 숨은 질적 구조를 드러내고,',
        '   AVM 실무의 설명 책무를 뒷받침하는 XAI 프레임워크 구축',
    ]),
    6: ('2절', '서 론', '연구의 범위', [
        '[공간적 범위]',
        '• 서울특별시 25개 자치구 · 215개 행정동',
        '',
        '[시간적 범위]',
        '• 2019년 1월 ~ 2025년 12월 (84개월, 7년)',
        '• COVID-19·금리인상·인하 전환의 거시 국면 모두 포함',
        '',
        '[분석 대상]',
        '• 아파트 매매 실거래 391,826건',
        '',
        '[데이터 출처]',
        '• 국토교통부 실거래가 (data.go.kr)',
        '• 서울열린데이터광장 — 지하철·백화점·CCTV·공원·도서관·어린이집',
        '• NEIS 교육정보 — 학교·학원',
        '• 한국은행 ECOS — 금리·CPI·M2',
    ]),
    7: ('2절', '서 론', '연구의 방법', [
        '[1단계] 데이터 수집 및 전처리',
        '• 법정동 → 행정동 매핑 (Nominatim 지오코딩 + GeoJSON 공간조인)',
        '• Y = log(거래금액 / 전용면적) 파생변수 계산',
        '',
        '[2단계] 3가지 모형 적합',
        '• OLS (베이스라인) · Random Forest · XGBoost',
        '• 3가지 분할: 무작위 / Group(단지) / 시간순',
        '',
        '[3단계] SHAP 해석',
        '• TreeSHAP (Lundberg et al., 2020), 테스트 5,000건 샘플',
        '• 글로벌 중요도 + Dependence Plot + Force Plot',
        '',
        '[4단계] 지역 별도 모형 및 강건성 점검',
        '• 전체·강남3구·비강남 22개구 독립 적합',
        '• Ablation · Moran\'s I · 시기별 3구간 보조 분석',
    ]),
    9: ('1절', '이론적 배경', '헤도닉 가격모형과 단위면적 정규화', [
        '[헤도닉 이론의 기초]',
        '• Lancaster (1966): 소비자 효용 = 재화의 속성 묶음에서 파생',
        '• Rosen (1974): 이질적 재화의 암묵적 가격 P = f(z₁, z₂, …, zₙ)',
        '• 전통 OLS 추정: P = β₀ + β₁z₁ + … + βₙzₙ + ε',
        '',
        '[OLS 헤도닉 추정의 3대 한계]',
        '• 선형 가정 → 비선형 가격 변동 미포착 (예: 재건축 U자형)',
        '• 변수 간 상호작용 효과 수동 지정 필요',
        '• 다중공선성에 취약 (거시 변수 간 VIF > 100)',
        '',
        '[본 연구의 설계: log(㎡당 가격)]',
        '• Y = ln(거래금액 / 전용면적) → 규모효과를 분모에 흡수',
        '• Malpezzi(2003), Sirmans et al.(2005) 표준 log-price 전통',
        '• 준 탄력성 해석: β → (e^β − 1)·100% / 단위 변화',
    ]),
    10: ('2절', '이론적 배경', '머신러닝 기반 예측 모형', [
        '[트리 기반 앙상블]',
        '• Random Forest (Breiman, 2001): 배깅 기반 독립 트리 평균',
        '• XGBoost (Chen & Guestrin, 2016): 부스팅 + L1/L2 정규화',
        '• LightGBM, CatBoost: 최신 부스팅 변형',
        '',
        '[본 연구 XGBoost 하이퍼파라미터]',
        '• max_depth = 8, learning_rate = 0.1',
        '• min_child_weight = 5, reg_alpha = 0.1, reg_lambda = 1.0',
        '• subsample = 0.8, colsample_bytree = 0.8, tree_method = hist',
        '• 최대 2,000 라운드 + early stopping patience = 50',
        '',
        '[Random Forest 설정]',
        '• n_estimators = 200, max_depth = 15, min_samples_leaf = 5',
        '• 50,000건 표본 + 3-Fold CV로 사전 탐색 후 채택',
    ]),
    11: ('3절', '이론적 배경', 'XAI와 SHAP (SHapley Additive exPlanations)', [
        '[SHAP의 이론적 기초]',
        '• Lundberg & Lee (2017): 게임이론 Shapley value를 ML 해석에 적용',
        '• 각 특성이 개별 예측에 미치는 기여도를 이론적으로 일관되게 분해',
        '• 블랙박스 모형의 해석 가능성 확보',
        '',
        '[TreeSHAP (Lundberg et al., 2020)]',
        '• 트리 기반 모형 전용 정확 Shapley 계산 알고리즘',
        '• 본 연구: 테스트 5,000건 샘플에 대해 SHAP값 산출',
        '',
        '[본 연구 활용 3계층]',
        '• 글로벌 중요도: mean(|SHAP_log|) 기준 18개 변수 서열',
        '• Dependence Plot: 변수값 vs SHAP의 비선형 관계',
        '• Force/Waterfall: 개별 거래 단위 기여 분해 (방법론 기반)',
        '',
        '[역변환 해석]',
        '→ exp(SHAP_log) ≈ 1 + SHAP_log → ±x% 단위가격 변화',
    ]),
    12: ('4절', '이론적 배경', '선행연구 검토 및 본 연구의 위치', [
        '[국내 ML 예측 연구]',
        '• 김이환 외(2022), 김학현 외(2023), 이선구(2025): 예측 성능 비교 중심',
        '• 조보근 외(2020): LIME 해석력 검증',
        '• 진수정(2024): 그래프 신경망 서울 아파트',
        '',
        '[해외 XAI 부동산 적용]',
        '• Neves et al.(2024): 리스본 스마트시티 SHAP',
        '• Mora-García et al.(2022): COVID-19 시기 ML+DL',
        '• Chun et al.(2025), Kim et al.(2025), An et al.(2025): 한국 초기 적용',
        '',
        '[본 연구의 차별성]',
        '• 단위면적 정규화 (규모효과 분리) × 행정동 × 지역 별도 모형 (3축 결합)',
        '• 3가지 분할 동시 보고 · Ablation · Moran\'s I 강건성',
        '• 국내 XAI 부동산 연구에서 최초 적용',
    ]),
    13: ('4절', '이론적 배경', '선행연구 상세 1 — 전통 헤도닉·공간계량', [
        '[국내 헤도닉·공간 연구]',
        '• 김우성 외(2019): 강남 아파트 헤도닉 구조 변화',
        '• 신광문·이재수(2019): 공간 헤도닉 소형 임대료',
        '• 장희선 외(2021): 공간 헤도닉 서울 공원 서비스 편익',
        '• Kim et al.(2022): 서울 주택 공간회귀 공간의존성',
        '',
        '[국내 거시·정책 연구]',
        '• 배종찬·정재호(2021): VEC 거시+부동산정책',
        '• 노산하·신진호(2021): 검색빈도 주택가격 인하효과',
        '• 정진오·정재호(2023): 조세·금융정책 차별효과 VAR',
        '• 김명진·서원석(2023): 코로나19 소매시설 시계열',
    ]),
    14: ('4절', '이론적 배경', '선행연구 상세 2 — 머신러닝 예측', [
        '[국내]',
        '• 배성완·유정석(2018): SVM/RF/GBT/DNN/LSTM 초기 비교',
        '• 김이환 외(2022): ML 아파트 매매가격지수',
        '• 김학현 외(2023): DNN·XGBoost·CatBoost 비교',
        '• 조보근 외(2020): 지역별 예측모델 + LIME',
        '• 이선구(2025): XGBoost AVM 은평 다세대 86.4%',
        '',
        '[해외]',
        '• Limsombunchai(2004): 헤도닉 vs ANN 초기 비교',
        '• Čeh et al.(2018): RF vs 다중회귀 성능 실증',
        '• Choy & Ho(2023): ML 부동산 적용 문헌 리뷰',
        '• Mora-García et al.(2022): COVID-19 시기 ML+DL 비교',
    ]),
    15: ('4절', '이론적 배경', '선행연구 상세 3 — XAI·SHAP 적용', [
        '[XAI 이론 기반]',
        '• Ribeiro et al.(2016): LIME 원 프레임워크',
        '• Lundberg & Lee(2017): SHAP 이론',
        '• Lundberg et al.(2020): TreeSHAP 확장',
        '',
        '[해외 XAI 부동산 적용]',
        '• Neves et al.(2024): 리스본 스마트시티 SHAP',
        '• Tarasov & Dessoulavy-Śliwiński(2025): XAI 헤도닉',
        '• Ezil Sam Leni et al.(2025): XGBoost+SHAP+웹배포',
        '',
        '[한국 XAI 초기 적용]',
        '• Chun et al.(2025): 서울 주택 XAI+ML 초기 적용',
        '• Kim et al.(2025): 한국 대량감정평가 XAI',
        '• An et al.(2025): 헤도닉 vs ML 투명성·정확성',
        '',
        '→ 국내 XAI 적용은 초기 단계, 단위가격·지역별 분해는 부재',
    ]),
    16: ('4절', '이론적 배경', '본 연구의 위치 설정 (종합)', [
        '[방법론적 차별화]',
        '① 종속변수: log(㎡당 가격) → 규모효과 제거',
        '② 공간 단위: 215개 행정동 (기존 25개 자치구 대비 세분)',
        '③ 지역 별도 모형: 전체 · 강남3구 · 비강남 22개구 독립 적합',
        '④ 이중 해석: SHAP 글로벌 + Dependence + Force (기반)',
        '',
        '[실증 설계 차별화]',
        '• 3가지 분할 동시 보고 (Random / Group / Chronological)',
        '• Ablation 실험 (학원·어린이집 제거)',
        '• Moran\'s I 공간 자기상관 검정',
        '• 시기별 3구간 보조 분석',
        '',
        '[해석 관점 차별화]',
        '• 준 탄력성 해석(%) — 실무 이해 용이',
        '• AVM 설명 책무(explainability obligation) 연결',
    ]),
    17: ('4절', '이론적 배경', '선행연구 종합 비교표', [
        '[구분 기준] 종속변수 | 공간 단위 | 모형 | 해석',
        '',
        '• Čeh et al.(2018): 총가격 | Ljubljana | OLS vs RF | 성능만',
        '• Chun et al.(2025): 총가격 | 서울 구 | XGB+SHAP | 통합 SHAP',
        '• Kim et al.(2025): 총가격 | 한국 | XAI 앙상블 | 통합 SHAP',
        '• An et al.(2025): 총가격 | 서울 | 헤도닉 vs ML | 투명성',
        '• Neves et al.(2024): 총가격 | 리스본 | XGB+SHAP | 공원 거리',
        '',
        '[본 연구]',
        '• Y: log(㎡당 가격) | 서울 215 행정동 | OLS + RF + XGB + SHAP',
        '• 지역 별도 3모형 · 3가지 분할 · Ablation · Moran\'s I',
        '',
        '→ 단위가격 + 지역 별도 + 다각 강건성 결합은 국내 최초',
    ]),
    19: ('1절', '연구 설계', '연구의 흐름도', [
        '[1단계] 데이터 수집',
        '   국토교통부 · 서울데이터광장 · NEIS · 한국은행 ECOS',
        '',
        '[2단계] 전처리 및 파생',
        '   법정동 → 행정동 매핑 (Nominatim + GeoJSON 공간조인)',
        '   Y = log(거래금액 / 전용면적) 계산',
        '',
        '[3단계] 모형 적합 (OLS · RF · XGBoost)',
        '   3가지 분할: 무작위 / Group(단지) / 시간순',
        '',
        '[4단계] SHAP 해석',
        '   TreeSHAP 기반 글로벌 · Dependence · Force',
        '',
        '[5단계] 지역 별도 모형 · 강건성 점검',
        '   전체 · 강남3구 · 비강남 22개구 독립 적합',
        '   Ablation · Moran\'s I · 시기별 3구간',
    ]),
    20: ('2절', '연구 설계', '변수 설정 (종속 1 + 독립 18)', [
        '[종속변수]',
        '• log(단위면적당 가격) = ln(거래금액[만원] / 전용면적[㎡])',
        '',
        '[독립변수 18개]',
        '• 물리적 특성 (3): 전용면적(㎡), 층, 건물연령(년)',
        '• 입지 특성 (2): 강남구분(0/1 더미), 지하철역수(개)',
        '• 환경 특성 (9):',
        '   초·중·고등학교수 / CCTV수 / 백화점수* / 공원수 / 도서관수',
        '   학원수† · 어린이집수† (구 단위 프록시)',
        '• 거시경제 (4): 기준금리, CD금리, 소비자물가지수, M2',
        '',
        '* 백화점수 = LOCALDATA_072405에서 "백화점" 매장 수',
        '† 프록시: 구 단위 총량을 행정동 수로 균등 배분',
        '※ 시설 변수는 수집시점 stock을 전 기간 병합 (연도별 변동 미반영)',
    ]),
    21: ('2절', '연구 설계', '3가지 데이터 분할 방식', [
        '[① 무작위 분할 (Random 70/10/20)]',
        '• 학습 274,277 / 검증 39,183 / 테스트 78,366',
        '• 동일 단지·평형 반복거래 공존 가능 → "동일 단지 내 예측" 성능',
        '',
        '[② Group 분할 (법정동+아파트명 기준)]',
        '• GroupShuffleSplit, 단지 수 8,601, overlap = 0',
        '• 미경험 단지 일반화 성능 (가장 보수적)',
        '',
        '[③ 시간순 분할 (Chronological)]',
        '• 학습: 2019.01 ~ 2023.12 (250,544건)',
        '• 검증: 2024.01 ~ 2024.06 (27,813건)',
        '• 테스트: 2024.07 ~ 2025.12 (113,469건)',
        '• 미래 시점 + 시장 국면 이동 반영',
    ]),
    22: ('3절', '연구 설계', '모형 구성 및 평가 지표', [
        '[OLS 다중회귀 (베이스라인)]',
        '• log(P/A) = β₀ + Σβᵢxᵢ + ε',
        '• 준 탄력성 해석: (e^β − 1) × 100% / 단위',
        '',
        '[Random Forest / XGBoost]',
        '• 동일 하이퍼파라미터 (앞 장 참고)',
        '• 모든 분할 · 모든 지역 모형에 동일 설정 유지',
        '',
        '[평가 지표]',
        '• R²_log: log 스케일 결정계수 (모형 간 동일 척도 비교)',
        '• RMSE_log / MAE_log: log 스케일 오차',
        '• MAPE / Median APE: 원 스케일(만원/㎡) 역변환 상대 오차',
        '',
        '→ 단위가격 분포가 우측 skewed이므로 평균과 중앙값을 병기',
    ]),
    23: ('4절', '연구 설계', 'SHAP 분석 방법', [
        '[TreeSHAP (Lundberg et al., 2020)]',
        '• XGBoost 전용 정확 Shapley 계산',
        '• 테스트 데이터에서 5,000건 무작위 샘플링',
        '',
        '[3가지 분석 층위]',
        '• 글로벌 변수 중요도',
        '   mean(|SHAP_log|) 기준 18개 변수 서열 및 비중(%)',
        '',
        '• Dependence Plot (비선형 검증)',
        '   변수값 ↔ SHAP 산점도 → 임계구간·체감·반등 패턴 확인',
        '',
        '• Force / Waterfall Plot (개별 해석)',
        '   base value → 최종 예측까지 변수별 기여 분해',
        '',
        '[역변환]',
        '→ exp(SHAP_log) ≈ ±x% 단위가격 변화 (실무 친화적 해석)',
    ]),
    24: ('5절', '연구 설계', '강건성 점검 및 지역 별도 모형', [
        '[강건성 점검 4종]',
        '① 3가지 분할 비교 (누수·국면 이동 분해)',
        '② Ablation: 학원·어린이집 제거 시 ΔR² 측정',
        '③ Moran\'s I: 잔차 공간 자기상관 검정',
        '④ 시기별 3구간: 유동성·금리인상·금리인하',
        '',
        '[지역 별도 모형 설계 — 교수 피드백 반영]',
        '• 기존 접근: 통합 모형 + SHAP 지역별 재집계',
        '• 본 연구: 3개 지역 독립 적합 (강남더미 제거, 17 features)',
        '   — 전체 (n=391,826)',
        '   — 강남3구 (n=65,077)',
        '   — 비강남 22개구 (n=326,749)',
        '• 동일 하이퍼파라미터 · 동일 70/10/20 무작위 분할',
        '• 각 지역 내부에서 독립 SHAP 구조 추출',
    ]),
    27: ('1절', '실증 분석', '기술통계 및 지역별 분포', [
        '[전체 표본] n = 391,826건 (2019.1~2025.12, 215개 행정동)',
        '',
        '[종속변수: ㎡당 가격 (만원/㎡)]',
        '• 평균 1,337.7 / 중앙값 1,142.7 / 최대 10,586.7',
        '• 25%분위 844.4 / 75%분위 1,614.9',
        '• log 변환 후 평균 7.0747, 표준편차 0.4879',
        '',
        '[지역별 단위가격]',
        '• 강남3구 (n=65,077): 평균 2,207.5만원/㎡',
        '• 비강남 (n=326,749): 평균 1,164.5만원/㎡',
        '• 배율 1.90배 (총가격 2.20배 → 면적효과 일부 분해)',
        '',
        '[독립변수 분포]',
        '• 전용면적 평균 75.76㎡, 중앙값 79.47㎡',
        '• 건물연령 평균 19.85년',
        '• 기준금리 0.50% ~ 3.50% (팬데믹·긴축·인하 전환 포함)',
    ]),
    28: ('2절', '실증 분석', '상관관계 및 VIF 진단', [
        '[log(㎡당 가격)과의 상관계수 (양의 방향 상위)]',
        '• 강남구분 r = +0.481 (가장 강한 연관)',
        '• 백화점수 r = +0.396 (상업 집적)',
        '• M2 통화량 r = +0.353 (유동성)',
        '• 소비자물가지수 r = +0.331 (인플레이션)',
        '• 층 r = +0.187',
        '',
        '[약한 또는 음의 상관]',
        '• 전용면적 r = +0.053 (총가격 0.581 → 규모효과 분리 증거)',
        '• 건물연령 r = −0.082 (재건축 U자형 비선형 배후)',
        '• 어린이집수 r = −0.296 (비강남 주거밀집지 혼재)',
        '',
        '[VIF Top (다중공선성 진단)]',
        '• 기준금리 110.11, CD금리 96.44 (거시변수 간 극심한 공선성)',
        '• CPI 42.87, M2 28.38',
        '• 어린이집 5.99, 나머지 미시변수 모두 <5 (허용 수준)',
    ]),
    29: ('3절', '실증 분석', 'Ablation 및 Moran\'s I 공간 자기상관', [
        '[Ablation: 학원수 · 어린이집수 제거]',
        '• OLS       ΔR²_log = −0.0542 (선형 모형이 프록시에 의존)',
        '• Random Forest  ΔR² = −0.0062',
        '• XGBoost   ΔR²_log = −0.0004 (영향 사실상 없음)',
        '',
        '→ XGBoost는 두 프록시 없이도 동등 예측 구조 유지',
        '→ 구 단위 배분 한계가 전체 결과를 왜곡하지 않음',
        '',
        '[Moran\'s I 잔차 공간 자기상관 (같은 구 내 인접 가중)]',
        '• OLS      I = 0.3247, Z = 8.882, p < 0.001 (강한 양의 자기상관)',
        '• XGBoost  I = 0.0042, Z = 0.252, p = 0.880 (통계적 유의성 없음)',
        '',
        '→ XGBoost 비선형 조합이 공간 구조 상당 부분 흡수',
        '→ 정밀 경계 기반 가중 행렬 재검증은 후속 과제',
    ]),
    30: ('4절', '실증 분석', '모형별 · 분할별 예측 성능 비교', [
        '[분할 | OLS | Random Forest | XGBoost]  (R²_log / Median APE)',
        '',
        '• 무작위  0.5061 / 22.04%   0.8937 / 7.51%   0.9554 / 4.16%',
        '• Group   0.4618 / 22.35%   0.7414 / 13.12%  0.8005 / 11.72%',
        '• 시간순  0.3939 / 24.74%   0.6955 / 14.82%  0.7972 / 12.61%',
        '',
        '[핵심 해석]',
        '• 모든 분할에서 XGBoost > RF >> OLS 일관된 서열',
        '• 무작위 R²=0.9554: 동일 단지 내 예측 성능',
        '• Group / 시간순 R² ≈ 0.80: 미경험 단지·미래 시점 일반화',
        '• Group ≈ 시간순 → 단지 누수 > 국면 이동 효과',
        '',
        '[일반화 참조값]',
        '→ 미경험 단지 예측 R²≈0.80, Median APE≈12% (실무 AVM 기준선)',
    ]),
    31: ('5절', '실증 분석', 'SHAP 결과 및 지역 별도 비교', [
        '[전체 모형 SHAP Top 6 — 규모 정규화 후 재편]',
        '• 건물연령 12.43% (1위, 총가격 모형에서 3위였음)',
        '• 강남구분 11.74%',
        '• 어린이집수 11.53% ★ (총가격 8위 → 3위)',
        '• M2 통화량 10.77%',
        '• 학원수 9.61% ★ (총가격 7위 → 5위)',
        '• 전용면적 9.55% (총가격 1위 22.2% → 6위, 규모효과 분리)',
        '',
        '[지역 별도 모형 SHAP Top 3]',
        '• 강남3구 (R²=0.9356): 건물연령 16.5% - CPI 13.3% - 백화점 13.0%',
        '• 비강남 (R²=0.9421): 건물연령 14.8% - CPI 11.8% - 전용면적 11.2%',
        '  비강남 4~6위: 어린이집 10.6% · M2 10.1% · 학원 8.6%',
        '',
        '→ 강남 = 재건축·상업 주도  /  비강남 = 주거권·사교육 주도',
    ]),
    33: ('1절', '결론', '시사점', [
        '[학술적 시사점]',
        '• 국내 아파트 시장에 XAI 체계적 적용 범위 확장',
        '• 단위면적 정규화 → 가려진 질적·입지적 설명 신호 드러냄',
        '• 지역 별도 모형으로 공간 이질성 직접 실증 (통합 SHAP 재집계 대비)',
        '',
        '[정책 참고점]',
        '• 강남 vs 비강남: 단위가격 결정 구조가 질적으로 상이',
        '• 재건축 규제는 강남에, 주거·교육 인프라는 비강남에 차별적 효과',
        '• M2·CPI 거시 변수는 시장 국면 대리신호 (인과 근거 ✗)',
        '• 건물연령 U자형 패턴 → 재건축 정책 연령 구간별 차등 검토',
        '',
        '[실무적 시사점 (AVM)]',
        '• 단위가격 SHAP 분해 = 면적에 가려지지 않은 질적 프리미엄 투명화',
        '• 프롭테크·감정평가 실무 XAI 설명 책무 지원',
    ]),
    34: ('2절', '결론', '연구의 한계 및 향후 과제', [
        '[한계]',
        '• 시설 변수 stock 정보 병합 (연도별 변동 미반영) ← 교수 지적',
        '• 학원수·어린이집수 구 단위 균등 배분 프록시 한계',
        '• 질적 변수 (브랜드·조망·재건축 단계) 미반영',
        '• Moran\'s I 간이 가중 행렬 (정밀 경계 기반 미적용)',
        '• 공간계량 모형(SAR/SEM) 직접 비교 미수행',
        '• 강남3구 표본 규모 제약 (비강남의 1/5)',
        '',
        '[향후 과제]',
        '• 연도별 시설 패널 구축 (행안부·교육부 연간 스냅샷)',
        '• GIS 기반 최근접 거리(nearest distance) 변수',
        '• Queen/Rook 경계 기반 정밀 공간 가중 행렬 재검증',
        '• LightGBM · CatBoost · TabNet 추가 비교',
        '• LIME · ALE · Permutation Importance 교차 XAI 검증',
        '• 수도권 · 지방 대도시로 분석 범위 확장',
    ]),
}

for slide_num_1based, (section_num, section_label, title, lines) in CONTENTS.items():
    idx = slide_num_1based - 1
    if idx >= len(prs.slides):
        continue
    slide = prs.slides[idx]
    # 페이지 번호 계산 (content 슬라이드 1부터)
    # 실제 논문식 페이지 넘버: slide_num_1based (단순히 슬라이드 순서 사용)
    write_content_slide(slide, page_num=slide_num_1based,
                        section_num=section_num, section_label=section_label,
                        title=title, lines=lines,
                        font_size=15, title_size=28)

# SHAP bar plot은 slide 31에 추가
from pptx.util import Inches as I
s31 = prs.slides[30]
shap_bar = os.path.join(PLOTS, 'fig5_shap_bar.png')
has_pic = any(str(sh.shape_type) == 'PICTURE (13)' for sh in s31.shapes)
if not has_pic and os.path.exists(shap_bar):
    s31.shapes.add_picture(shap_bar, I(11.5), I(6.5), height=I(4.2))

# 참고문헌(Slide 35) — 본 연구 참고문헌으로 전면 교체
s = prs.slides[34]
# 기존 양재영 참고문헌 모두 제거
for sh in list(s.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip():
        remove_shape(sh)
# 제목
title_box = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(18), Inches(0.8))
tf = title_box.text_frame
add_p(tf, '참고문헌', size=32, bold=True, color=DARK, first=True)

# 좌측 국내문헌
left_box = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(9), Inches(9.5))
tf = left_box.text_frame; tf.word_wrap = True
add_p(tf, '[국내문헌]', size=14, bold=True, color=CYAN, first=True)
kor_refs = [
    '김명진, 서원석 (2023). 코로나19에 따른 소매유통시설의 아파트 가격 영향. 토지주택연구, 14(3), 37-53.',
    '김우성 외 (2019). 헤도닉 가격 모형을 통한 강남 지역 주거 선호 구조 변화. 부동산학보, 76, 137-150.',
    '김이환 외 (2022). 기계학습 방법론 아파트 매매가격지수 연구. 부동산분석, 8(3), 1-29.',
    '김학현 외 (2023). 딥러닝·머신러닝 아파트 실거래가 예측. 정보처리학회 논문지, 12(2), 59-76.',
    '노산하, 신진호 (2021). 검색빈도 종합부동산세·재산세 주택가격 효과. 국토연구, 110, 81-93.',
    '배성완, 유정석 (2018). ML·시계열 분석 부동산 가격지수 예측. 주택연구, 26(1), 107-133.',
    '배종찬, 정재호 (2021). 거시경제·부동산정책 서울 아파트가격 영향. 토지주택연구, 12(4), 41-59.',
    '신광문, 이재수 (2019). 공간 헤도닉 소형주택 임대료 결정. 부동산분석, 5(3), 49-66.',
    '이선구 (2025). XGBoost AVM 은평구 다세대 실증. 부동산분석, 11(2), 21-40.',
    '장희선 외 (2021). 공간헤도닉 서울 도시공원 편익측정. 국토계획, 56(5), 215-227.',
    '정진오, 정재호 (2023). 조세·금융정책 부동산가격 영향. 토지주택연구, 14(3), 55-75.',
    '조보근 외 (2020). ML 알고리즘 지역별 아파트 예측 + LIME. 정보시스템연구, 29(3), 119-144.',
    '진수정 (2024). SRGCNN 서울시 아파트 매매가격 예측. 서울대학교 석사논문.',
    '통계청 (2024). 가계금융복지조사 결과.',
]
for r in kor_refs:
    add_p(tf, '• ' + r, size=10, color=DARK)

# 우측 해외문헌
right_box = s.shapes.add_textbox(Inches(10), Inches(1.4), Inches(9.2), Inches(9.5))
tf = right_box.text_frame; tf.word_wrap = True
add_p(tf, '[해외문헌]', size=14, bold=True, color=CYAN, first=True)
intl_refs = [
    'An, S. et al. (2025). Hedonic vs ML housing appraisal. Financial Innovation, 11:141.',
    'Anselin, L. (1988). Spatial Econometrics. Kluwer.',
    'Breiman, L. (2001). Random forests. Machine Learning, 45(1), 5-32.',
    'Čeh, M. et al. (2018). RF vs regression apartment prices. ISPRS IJGI, 7(5), 168.',
    'Chen, T. & Guestrin, C. (2016). XGBoost. KDD\'16, 785-794.',
    'Choy, L. & Ho, W. (2023). ML in real estate research. Land, 12(4), 740.',
    'Chun, H. et al. (2025). Seoul XAI+ML housing price. KSII, 19(4), 1077-1096.',
    'Friedman, J. (2001). Gradient boosting machine. Annals of Stat., 29(5).',
    'Kim, J.-J. et al. (2022). Spatial regression multiplex houses. Sustainability, 14(12), 7116.',
    'Kim, W. et al. (2025). Explainable AI mass appraisal Korea. IJSPM, 29(5), 350-376.',
    'Lancaster, K. (1966). Consumer theory. J. Political Economy, 74(2).',
    'Limsombunchai, V. (2004). Hedonic vs ANN. NZARES Conf.',
    'Lundberg, S. & Lee, S. (2017). SHAP. NeurIPS 30, 4765-4774.',
    'Lundberg, S. et al. (2020). TreeSHAP. Nature MI, 2(1), 56-67.',
    'Mora-García, R. et al. (2022). ML COVID-19 housing. Land, 11(11), 2100.',
    'Neves, F. et al. (2024). Open data + XAI real estate. Applied Sciences, 14(5), 2209.',
    'Ribeiro, M. et al. (2016). LIME. KDD\'16, 1135-1144.',
    'Rosen, S. (1974). Hedonic prices. J. Political Economy, 82(1).',
    'Tarasov, S. & Dessoulavy-Ś. (2025). XAI hedonic pricing. REMV, 33(1), 22-34.',
]
for r in intl_refs:
    add_p(tf, '• ' + r, size=10, color=DARK)

fix_font_eastasian(s)

# 감사합니다(Slide 36) — 원본 유지 + 폰트 정규화
s = prs.slides[35]
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    r.font.name = FONT_KR
                    set_font_ea(r)

prs.save(DST)
print(f"생성 완료: {DST}")
print(f"파일 크기: {os.path.getsize(DST)/1024/1024:.1f} MB")
