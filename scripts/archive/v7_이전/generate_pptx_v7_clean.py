#!/usr/bin/env python3
"""v7 발표 PPT 깨끗한 재생성:
- 양재영 템플릿의 레이아웃(섹션 표지 + 페이지 번호 + 색상)을 유지
- 모든 신규 텍스트에 Pretendard 폰트 명시 → 한글 깨짐 방지
- 본문 슬라이드는 기존 텍스트를 전부 지우고 새 콘텐츠로 재작성
- 잔존 YJY 텍스트 제거 + 빈 placeholder 정리"""
import os, shutil, copy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = '/Users/seokmogu/Downloads/양재영_중간발표_1203.pptx'
DST = '/Users/seokmogu/project/thesis-seoul-apt-xai/paper/중간발표_v7_m2price.pptx'
PLOTS = '/Users/seokmogu/project/thesis-seoul-apt-xai/results/plots_v7_m2price'

# 템플릿 복사본 새로 만들기
shutil.copy(SRC, DST)
prs = Presentation(DST)

CYAN = RGBColor(0x08, 0xA5, 0xC1)
DARK = RGBColor(0x39, 0x39, 0x39)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_KR = 'Pretendard'          # 한글 본문
FONT_KR_BOLD = 'Pretendard'     # Bold 가중치는 run.font.bold로 제어
FONT_EN = 'Montserrat'          # 숫자·영문

# ========== 공통 유틸 ==========
def clear_text_frame(shape):
    """shape의 TextFrame 모든 paragraph 제거, 빈 상태로."""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)
    # 최소 1개 paragraph 필요
    etree.SubElement(txBody, qn('a:p'))
    return True

def add_line(tf, text, size=14, bold=False, color=DARK, font=FONT_KR, alignment=None, first=False):
    """TextFrame에 텍스트 라인 추가 (first=True면 기존 첫 paragraph 사용)."""
    if first:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    if alignment is not None:
        para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    try: run.font.color.rgb = color
    except: pass
    # East Asian font 명시 (한글 깨짐 방지)
    rPr = run._r.get_or_add_rPr()
    # eastAsia 폰트 설정
    for existing in rPr.findall(qn('a:ea')):
        rPr.remove(existing)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    for existing in rPr.findall(qn('a:latin')):
        rPr.remove(existing)
    latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', font)
    return run

def remove_shape(shape):
    try:
        shape._element.getparent().remove(shape._element)
    except: pass

def replace_content(slide, title, lines):
    """슬라이드의 큰 텍스트 박스를 찾아 title + lines로 교체.
    모든 텍스트를 Pretendard로 강제 적용."""
    # 가장 큰 텍스트 박스 찾기 (본문 영역)
    text_shapes = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = sh.text_frame.text
            # 페이지 번호(숫자만) / 섹션 넘버(01~05) / 작은 라벨은 건드리지 않음
            if len(txt.strip()) <= 3 and txt.strip().replace('절','').isdigit():
                continue
            if txt.strip() in ['1절','2절','3절','4절','5절','01','02','03','04','05']:
                continue
            w, h = sh.width or 0, sh.height or 0
            text_shapes.append((w*h, sh, txt))
    if not text_shapes:
        return False
    text_shapes.sort(reverse=True, key=lambda x: x[0])
    main = text_shapes[0][1]
    # 지우고 새로 쓰기
    clear_text_frame(main)
    tf = main.text_frame
    tf.word_wrap = True
    # 제목
    if title:
        add_line(tf, title, size=28, bold=True, color=CYAN, font=FONT_KR, first=True)
        for ln in lines:
            if ln.strip() == '':
                add_line(tf, ' ', size=8, color=DARK, font=FONT_KR)
            else:
                size = 15
                bold = False
                if ln.startswith('[') and ln.endswith(']'):
                    size, bold = 16, True
                elif ln.startswith('①') or ln.startswith('②') or ln.startswith('③'):
                    size = 15
                elif ln.startswith('•') or ln.startswith('  '):
                    size = 14
                add_line(tf, ln, size=size, bold=bold, color=DARK, font=FONT_KR)
    else:
        for i, ln in enumerate(lines):
            if i == 0:
                add_line(tf, ln, size=15, color=DARK, font=FONT_KR, first=True)
            else:
                add_line(tf, ln, size=15, color=DARK, font=FONT_KR)
    return True

def remove_shapes_matching(slide, patterns, max_len=500):
    """지정 패턴 포함 shape 제거"""
    to_rm = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            for p in patterns:
                if p in t and (max_len is None or len(t) < max_len):
                    to_rm.append(sh); break
    for sh in to_rm:
        remove_shape(sh)

# ========== 슬라이드별 콘텐츠 ==========

# === Slide 1 (표지) ===
s = prs.slides[0]
# 기존 텍스트 전부 제거
for sh in list(s.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip():
        remove_shape(sh)
# 신규 표지
W = prs.slide_width
H = prs.slide_height
from pptx.util import Inches
# 제목 상단
title_box = s.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(17), Inches(2.8))
tf = title_box.text_frame
tf.word_wrap = True
add_line(tf, 'XGBoost와 SHAP을 활용한', size=36, bold=True, color=WHITE,
         font=FONT_KR, alignment=PP_ALIGN.CENTER, first=True)
add_line(tf, '서울시 아파트 단위면적당 매매가격의 설명 패턴 분석', size=44, bold=True,
         color=CYAN, font=FONT_KR, alignment=PP_ALIGN.CENTER)

# 부제
sub_box = s.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(17), Inches(0.8))
tf = sub_box.text_frame
add_line(tf, '단위면적 정규화 · 행정동 분석 · 지역 별도 모형', size=22,
         color=WHITE, font=FONT_KR, alignment=PP_ALIGN.CENTER, first=True)

# 소속/발표자/일자
info_box = s.shapes.add_textbox(Inches(2), Inches(8.5), Inches(16), Inches(1.8))
tf = info_box.text_frame
add_line(tf, '한양대학교 부동산융합대학원 · 도시부동산정책전공', size=18,
         color=WHITE, font=FONT_KR, alignment=PP_ALIGN.CENTER, first=True)
add_line(tf, '석사학위 중간발표 | 박 현 근', size=24, bold=True,
         color=WHITE, font=FONT_KR, alignment=PP_ALIGN.CENTER)
add_line(tf, '2026. 04.', size=16, color=WHITE, font=FONT_KR, alignment=PP_ALIGN.CENTER)

# === Slide 2 (목차) — 기존 숫자·섹션 구조 유지, 라벨만 수정 ===
s = prs.slides[1]
# 기존 '서론/이론/분석의/분석결과/결론' 라벨 교체 및 "목차" 타이틀 유지
replace_map = {
    '서론': '서 론',
    '이론': '이론적 배경 및 선행연구',
    '분석의': '연구 설계 및 방법',
    '분석결과': '실증 분석 결과',
    '결론': '결론 및 시사점',
}
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                for old, new in replace_map.items():
                    if r.text.strip() == old:
                        r.text = new
                        r.font.name = FONT_KR
# 목차 제목(큰 "목차" 글자)에 폰트 강제
for sh in s.shapes:
    if sh.has_text_frame:
        t = sh.text_frame.text.strip()
        if t == '목차':
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT_KR

# === Slide 3 (섹션 표지 01 서론) ===
s = prs.slides[2]
# 큰 제목 '서론' 유지, 작은 라벨만 정리
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() == '연구의':
                    r.text = '서 론'
                    r.font.name = FONT_KR
                elif r.text.strip() == '서론':
                    r.font.name = FONT_KR

# === Slide 4: 연구의 배경 ===
s = prs.slides[3]
# 기존 표/그림 박스 제거
remove_shapes_matching(s, ['<표 1>', '<표 2>', '1인 가구 비율', '거처 종류'])
replace_content(s, '연구의 배경', [
    '• 한국 가계 자산의 75.2%가 실물자산, 대부분 부동산 (통계청, 2024)',
    '• 서울 아파트 = 자산 축적·주거 안정·거시경제 파급효과의 핵심 매개',
    '',
    '• 전통적 OLS 헤도닉 모형: 선형 가정, 비선형·상호작용 포착 한계',
    '• ML(RF·XGBoost): 예측 성능 우수하나 "블랙박스" 문제',
    '• XAI(SHAP, Lundberg & Lee 2017): 설명 가능성 확보 수단',
    '',
    '• 국내 아파트 시장에 XAI를 체계적으로 적용한 연구는 부족',
])

# === Slide 5: 연구 필요성과 목적 ===
s = prs.slides[4]
remove_shapes_matching(s, ['전월세', '주거사다리', '효과성', '주거 안정', '주거 정책'])
replace_content(s, '연구의 필요성 및 목적', [
    '[연구 동기: 왜 이 주제인가]',
    '• AVM(자동감정평가) 실무 확산 ↔ 설명 프레임워크 국내 부재',
    '• 기존 한국 ML 연구의 3공백:',
    '   ① 자치구(25개) 단위 집계  ② 단일 모형 예측 비교  ③ 통합 SHAP만',
    '• 총가격 모형: 면적이 SHAP 상위 독점 → 질적 신호 가려짐',
    '',
    '[본 연구의 3가지 차별성]',
    '① 단위면적당 가격 정규화 log(거래금액/전용면적) — 규모효과 분리',
    '② 행정동(215개) 단위 세분화 — 기존 자치구 대비',
    '③ 지역 별도 모형 (전체·강남3구·비강남22개구) 직접 비교',
])

# === Slide 6: 연구 범위 ===
s = prs.slides[5]
replace_content(s, '연구의 범위', [
    '[공간] 서울특별시 25개 자치구 · 215개 행정동',
    '[시간] 2019년 1월 ~ 2025년 12월 (84개월)',
    '[대상] 아파트 매매 실거래 391,826건',
    '',
    '[데이터 출처]',
    '• 국토교통부 실거래가 (data.go.kr)',
    '• 서울열린데이터광장 (data.seoul.go.kr) — 지하철·백화점·CCTV·공원·도서관·어린이집',
    '• NEIS 교육정보 (open.neis.go.kr) — 학교·학원',
    '• 한국은행 ECOS (ecos.bok.or.kr) — 금리·CPI·M2',
    '',
    '[종속변수] log(거래금액[만원] ÷ 전용면적[㎡])',
    '[독립변수] 18개 (물리 3 + 입지 2 + 환경 9 + 거시 4)',
])

# === Slide 7: 연구 방법 ===
s = prs.slides[6]
remove_shapes_matching(s, ['청년안심주택 정의', '공공임대', '한계 및'])
replace_content(s, '연구의 방법', [
    '[1단계] 데이터 수집·전처리',
    '• 법정동 → 행정동 매핑 (Nominatim 지오코딩 + GeoJSON 공간조인)',
    '• 거시경제 월별 데이터 병합, 시설 stock 정보 병합',
    '',
    '[2단계] 모형 구축 (동일 하이퍼파라미터)',
    '• OLS (베이스라인) / Random Forest / XGBoost',
    '• 3가지 분할 방식: 무작위 / Group(단지) / 시간순',
    '',
    '[3단계] SHAP 분석',
    '• TreeSHAP (Lundberg et al., 2020), 테스트 5,000건 샘플',
    '• 글로벌 중요도 + Dependence Plot + Force Plot 기반 구조',
    '',
    '[4단계] 지역 별도 모형 및 강건성 점검',
    '• 전체/강남/비강남 3모형 독립 적합',
    '• Ablation · Moran\'s I · 시기별 3구간 보조 분석',
])

# === Slide 8 (섹션 표지 02) ===
s = prs.slides[7]
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if '이론' in r.text or '선행연구' in r.text:
                    r.text = '이론적 배경 및 선행연구'
                    r.font.name = FONT_KR

# === Slide 9: 헤도닉 가격모형 ===
s = prs.slides[8]
replace_content(s, '헤도닉 가격모형 (Hedonic Price Model)', [
    '[이론적 기초]',
    '• Lancaster (1966): 소비자 효용 = 재화의 속성 묶음에서 파생',
    '• Rosen (1974): 이질적 재화의 암묵적 가격(implicit price)',
    '• P = f(z₁, z₂, …, zₙ), ∂P/∂zᵢ = pᵢ (한계 속성가격)',
    '',
    '[전통 OLS 헤도닉 추정의 한계]',
    '• 선형 가정 → 비선형 가격 변동 패턴 미포착 (예: 재건축 U자형)',
    '• 변수 간 상호작용 효과 수동 지정 필요',
    '• 다중공선성에 취약 (특히 거시경제 변수 간)',
    '',
    '[본 연구의 설계]',
    '• log(거래금액/전용면적)을 종속변수로 설정 → 규모효과 분모 제거',
    '• Malpezzi(2003), Sirmans et al.(2005) 표준 log-price 전통에 부합',
    '• 연속 독립변수 β → (e^β − 1)·100% 준 탄력성 해석',
])

# === Slide 10: 머신러닝 기반 예측 ===
s = prs.slides[9]
replace_content(s, '머신러닝 기반 부동산 가격 예측', [
    '[트리 기반 앙상블]',
    '• Random Forest (Breiman, 2001): 배깅 기반 독립 트리 앙상블',
    '• XGBoost (Chen & Guestrin, 2016): 그래디언트 부스팅 + 정규화',
    '• LightGBM, CatBoost: 최신 부스팅 변형',
    '',
    '[본 연구 XGBoost 하이퍼파라미터]',
    '• max_depth = 8, learning_rate = 0.1, min_child_weight = 5',
    '• reg_alpha = 0.1, reg_lambda = 1.0, tree_method = hist',
    '• subsample = 0.8, colsample_bytree = 0.8',
    '• 최대 2,000 라운드 + early stopping patience = 50',
    '',
    '[Random Forest]',
    '• n_estimators = 200, max_depth = 15, min_samples_leaf = 5',
    '• 사전 탐색: 50,000건 표본 3-Fold CV로 조합 선택',
])

# === Slide 11: XAI (SHAP) ===
s = prs.slides[10]
replace_content(s, 'XAI와 SHAP (SHapley Additive exPlanations)', [
    '[SHAP 개요 (Lundberg & Lee, 2017)]',
    '• 게임이론의 Shapley value를 ML 예측 해석에 적용',
    '• 각 특성의 개별 예측 기여도를 이론적으로 일관되게 분해',
    '• 블랙박스 모형의 해석 가능성 확보',
    '',
    '[TreeSHAP (Lundberg et al., 2020)]',
    '• 트리 기반 모형 전용 정확 계산 알고리즘',
    '• 본 연구: 테스트 샘플 5,000건 대상 SHAP값 산출',
    '',
    '[본 연구 활용 3단계]',
    '• 글로벌 변수 중요도: 평균 |SHAP| 기준 서열',
    '• Dependence Plot: 변수값 ↔ SHAP 비선형 관계',
    '• Force/Waterfall Plot: 개별 거래 단위 기여 분해 (방법론적 토대)',
])

# === Slide 12: 선행연구 종합 ===
s = prs.slides[11]
replace_content(s, '선행연구 검토 및 본 연구 차별성', [
    '[국내 선행연구의 3가지 공백]',
    '• 김이환 외(2022), 김학현 외(2023), 이선구(2025): 예측 성능 중심',
    '• 배성완·유정석(2018), 조보근 외(2020): ML 초기 적용',
    '• 자치구(25개) 단위 집계에 머무름',
    '',
    '[해외 XAI 적용 연구]',
    '• Neves et al.(2024): 리스본 스마트시티 SHAP 분석',
    '• Mora-García et al.(2022): COVID-19 시기 ML+DL 비교',
    '• Chun et al.(2025), Kim et al.(2025), An et al.(2025): 한국 초기 적용',
    '',
    '[본 연구 차별성]',
    '① 단위면적 정규화 (규모효과 분리) + 행정동 + 지역 별도 모형 (3축)',
    '② 3가지 분할 동시 보고 (무작위/Group/시간순)',
    '③ Ablation + Moran\'s I 강건성 점검',
])

# === Slide 13: 국내·해외 선행연구 상세 (헤도닉·거시) ===
s = prs.slides[12]
replace_content(s, '선행연구 상세 1: 전통 헤도닉·공간계량', [
    '[국내 헤도닉/공간 연구]',
    '• 김우성 외(2019): 강남 아파트 헤도닉 구조 변화',
    '• 신광문·이재수(2019): 공간 헤도닉 소형 임대료',
    '• 장희선 외(2021): 공간 헤도닉 서울 공원 서비스 편익',
    '• Kim et al.(2022): 서울 주택 공간회귀 공간의존성',
    '',
    '[국내 거시·정책 연구]',
    '• 배종찬·정재호(2021): VEC 거시+부동산정책',
    '• 노산하·신진호(2021): 검색빈도 주택가격 인하효과',
    '• 정진오·정재호(2023): 조세·금융정책 차별효과 VAR',
    '• 김명진·서원석(2023): 코로나19 소매시설 시계열',
])

# === Slide 14: ML 선행연구 ===
s = prs.slides[13]
replace_content(s, '선행연구 상세 2: 머신러닝 예측 연구', [
    '[국내]',
    '• 배성완·유정석(2018): SVM/RF/GBT/DNN/LSTM 비교',
    '• 김이환 외(2022): ML 아파트 매매가격지수',
    '• 김학현 외(2023): DNN·XGBoost·CatBoost 비교',
    '• 조보근 외(2020): 지역별 아파트 예측모델 + LIME',
    '• 이선구(2025): XGBoost AVM 은평 다세대 86.4%',
    '• 진수정(2024): 그래프 신경망(SRGCNN) 서울 아파트',
    '',
    '[해외]',
    '• Limsombunchai(2004): 헤도닉 vs ANN 초기 비교',
    '• Čeh et al.(2018): RF vs 다중회귀 성능 실증',
    '• Choy & Ho(2023): ML 부동산 적용 문헌 리뷰',
    '• Mora-García et al.(2022): COVID-19 시기 ML+DL',
])

# === Slide 15: XAI 선행연구 ===
s = prs.slides[14]
replace_content(s, '선행연구 상세 3: XAI·SHAP 적용', [
    '[해외 XAI/SHAP 부동산 적용]',
    '• Ribeiro et al.(2016): LIME 원 프레임워크',
    '• Lundberg & Lee(2017): SHAP 이론',
    '• Lundberg et al.(2020): TreeSHAP 확장',
    '• Neves et al.(2024): 리스본 스마트시티 SHAP',
    '• Tarasov & Dessoulavy-Śliwiński(2025): XAI 헤도닉',
    '• Ezil Sam Leni et al.(2025): XGBoost+SHAP 앙상블',
    '',
    '[한국 XAI 초기 적용]',
    '• Chun et al.(2025): 서울 주택 XAI+ML 초기 적용',
    '• Kim et al.(2025): 한국 주거 ML 대량감정평가',
    '• An et al.(2025): 헤도닉 vs ML 투명성·정확성',
    '',
    '→ 국내 XAI 적용은 여전히 초기 단계, 단위가격·지역별 분해 부재',
])

# === Slide 16: 본 연구 위치 설정 ===
s = prs.slides[15]
replace_content(s, '본 연구의 위치 설정', [
    '[방법론적 차별화]',
    '① 종속변수: log(㎡당 가격) → 규모효과 제거',
    '② 공간 단위: 215개 행정동 (기존 25개 자치구 대비 세분화)',
    '③ 지역 별도 모형: 강남3구 vs 비강남 22개구 vs 전체 독립 적합',
    '④ 이중 해석: SHAP 글로벌 + Dependence + 지역별',
    '',
    '[실증 설계 차별화]',
    '• 3가지 분할 동시 보고 (random / group / chronological)',
    '• Ablation 실험 (학원·어린이집 제거)',
    '• Moran\'s I 공간 자기상관 검정',
    '• 시기별 3구간 (유동성/금리인상/금리인하) 보조 분석',
    '',
    '[해석 관점 차별화]',
    '• 준 탄력성 해석 (%/단위) — 실무 이해 용이',
    '• AVM 설명 책무(explainability obligation) 연결',
])

# === Slide 17: 선행연구 요약표 (빈 슬라이드로 정리) ===
s = prs.slides[16]
replace_content(s, '선행연구 종합 비교', [
    '[구분 기준] 종속변수 | 공간 단위 | 모형 | 해석 수단',
    '',
    '• Čeh et al.(2018): 총가격 | Ljubljana 시 | OLS vs RF | 성능만',
    '• Chun et al.(2025): 총가격 | 서울 구 | XGB+SHAP | 통합 SHAP만',
    '• Kim et al.(2025): 총가격 | 한국 | XAI 앙상블 | 통합 SHAP',
    '• An et al.(2025): 총가격 | 서울 | 헤도닉 vs ML | 투명성 비교',
    '• Neves et al.(2024): 총가격 | 리스본 | XGB+SHAP | 공원 거리',
    '',
    '[본 연구]',
    '• log(㎡당 가격) | 서울 215 행정동 | OLS+RF+XGB+SHAP',
    '• 지역 별도 3모형 · 3가지 분할 · Ablation · Moran\'s I',
    '',
    '→ 단위가격 + 지역 별도 + 다각 강건성의 결합은 국내 최초 적용',
])

# === Slide 18 (섹션 표지 03) ===
s = prs.slides[17]
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if '분석의' in r.text or '분석' in r.text.strip():
                    if len(r.text.strip()) < 10:
                        r.text = '연구 설계 및 방법'
                        r.font.name = FONT_KR

# === Slide 19: 연구 흐름 ===
s = prs.slides[18]
remove_shapes_matching(s, ['[ 1단계 ]', '[ 2단계 ]', '[ 3단계 ]', '[ 4단계 ]', '[ 5단계 ]'])
replace_content(s, '연구의 흐름도', [
    '[1단계] 데이터 수집',
    '  국토교통부 · 서울데이터광장 · NEIS · 한국은행 ECOS',
    '',
    '[2단계] 전처리 및 파생',
    '  법정동→행정동 매핑 (Nominatim + GeoJSON)',
    '  Y = log(거래금액 / 전용면적) 계산',
    '',
    '[3단계] 모형 적합 (OLS · RF · XGBoost)',
    '  3가지 분할: 무작위 / Group(단지) / 시간순',
    '',
    '[4단계] SHAP 해석',
    '  TreeSHAP 기반 글로벌·Dependence·Force',
    '',
    '[5단계] 지역 별도 모형 · 강건성 점검',
    '  전체·강남3구·비강남 22개구 독립 적합',
    '  Ablation · Moran\'s I · 시기별 3구간',
])

# === Slide 20: 변수 설정 ===
s = prs.slides[19]
replace_content(s, '변수 설정 (18개 독립변수)', [
    '[물리적 특성 3개]',
    '• 전용면적 (㎡), 층, 건물연령 (년)',
    '',
    '[입지 특성 2개]',
    '• 강남구분 (0/1 더미), 지하철역수 (개)',
    '',
    '[환경 특성 9개]',
    '• 초·중·고등학교수, CCTV수, 백화점수*, 공원수, 도서관수',
    '• 학원수† · 어린이집수† (구 단위 프록시)',
    '',
    '[거시경제 4개]',
    '• 기준금리, CD금리, 소비자물가지수, M2 (월별 stock)',
    '',
    '* 백화점수: LOCALDATA_072405에서 "백화점" 매장 수',
    '† 프록시: 구 단위 총량을 행정동 수로 균등 배분',
    '※ 시설 변수는 수집시점 stock을 거래 전 기간에 병합 (한계)',
])

# === Slide 21: 분할 방식 ===
s = prs.slides[20]
replace_content(s, '3가지 데이터 분할 방식', [
    '[① 무작위 분할 (Random 70/10/20)]',
    '• 학습 274,277 / 검증 39,183 / 테스트 78,366',
    '• 동일 단지·동일 평형 반복거래 가능 → "동일 단지 내 예측" 성능',
    '',
    '[② Group 분할 (법정동+아파트명 기준)]',
    '• GroupShuffleSplit 단지 수 8,601, overlap = 0',
    '• 미경험 단지 일반화 성능 측정 (가장 보수적)',
    '',
    '[③ 시간순 분할 (Chronological)]',
    '• 학습: 2019.01 ~ 2023.12 (250,544건)',
    '• 검증: 2024.01 ~ 2024.06 (27,813건)',
    '• 테스트: 2024.07 ~ 2025.12 (113,469건)',
    '• 미래 시점 예측 + 시장 국면 이동 반영',
])

# === Slide 22: 모형 구성 ===
s = prs.slides[21]
replace_content(s, '모형 구성 및 평가 지표', [
    '[OLS 다중회귀 (베이스라인)]',
    '• log(P/A) = β₀ + Σβᵢxᵢ + ε',
    '• 준 탄력성 해석: (e^β − 1) × 100% / 단위',
    '',
    '[Random Forest]',
    '• n_estimators = 200, max_depth = 15, min_samples_leaf = 5',
    '',
    '[XGBoost (대표 해석 모형)]',
    '• max_depth = 8, learning_rate = 0.1, 최대 2,000 라운드',
    '• early_stopping patience = 50 (미발동, 1,999 라운드 종료)',
    '',
    '[평가 지표]',
    '• R²_log (log 스케일 결정계수) — 모형 간 동일 척도 비교',
    '• RMSE_log, MAE_log — log 스케일 오차',
    '• MAPE · Median APE — 원 스케일(만원/㎡) 역변환 상대 오차',
])

# === Slide 23: SHAP 분석 방법 ===
s = prs.slides[22]
replace_content(s, 'SHAP 분석 방법', [
    '[TreeSHAP (Lundberg et al., 2020)]',
    '• XGBoost 전용 정확 Shapley value 계산 알고리즘',
    '• 샘플: 테스트 데이터에서 5,000건 무작위 추출',
    '',
    '[3가지 SHAP 분석 층위]',
    '• 글로벌 변수 중요도',
    '   mean(|SHAP_log|) 기준 18개 변수 서열 및 비중(%)',
    '',
    '• Dependence Plot (비선형 검증)',
    '   변수값 vs SHAP 기여도의 산점도 → 임계구간·체감·반등 확인',
    '',
    '• Force/Waterfall Plot (개별 해석)',
    '   특정 거래의 base value → 최종 예측 변수별 기여 분해',
    '',
    '[역변환 해석]',
    '• exp(SHAP_log) ≈ 1 + SHAP_log → 약 ±x% 단위가격 변화',
])

# === Slide 24: 강건성 점검 ===
s = prs.slides[23]
replace_content(s, '강건성 점검 설계', [
    '[① 3가지 분할 비교]',
    '• 무작위 ↔ Group ↔ 시간순: 성능 격차의 구성 요인 이원 분해',
    '',
    '[② Ablation 실험]',
    '• 학원수·어린이집수 제거 시 ΔR² 측정',
    '• 구 단위 프록시의 예측력 의존도 정량 평가',
    '',
    '[③ Moran\'s I 공간 자기상관 검정]',
    '• 행정동 잔차 평균에 대해 "같은 구 내 인접" 가중 행렬 적용',
    '• 499회 순열 검정 기반 p값 산출',
    '• OLS vs XGBoost 잔차 구조 비교',
    '',
    '[④ 시기별 3구간 (보조)]',
    '• 유동성장세(2019-2021) · 금리인상(2022-2023) · 금리인하(2024-2025)',
    '• 각 시기별 XGBoost 독립 적합 → 핵심 변수 안정성 점검',
])

# === Slide 25: 지역 별도 모형 설계 ===
s = prs.slides[24]
replace_content(s, '지역 별도 모형 설계 (교수 피드백 반영)', [
    '[기존 접근: 통합 모형 + SHAP 지역별 재집계]',
    '• 통합 XGBoost (강남구분 더미 포함) 학습',
    '• 예측·SHAP값을 지역별로 분리 집계',
    '• 한계: 통합 모형의 공통 구조에 강남 신호가 희석',
    '',
    '[본 연구: 3지역 독립 적합]',
    '• 전체 (17 features, 강남더미 제거): n = 391,826',
    '• 강남3구 (강남·서초·송파): n = 65,077',
    '• 비강남 22개구: n = 326,749',
    '',
    '[설계 특징]',
    '• 동일 하이퍼파라미터, 동일 70/10/20 무작위 분할',
    '• 각 지역 내부에서 독립적으로 SHAP 구조 추출',
    '• 연도별 분해 대신 공간 이질성 실증에 초점',
])

# === Slide 26 (섹션 표지 04) ===
s = prs.slides[25]
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if '분석결과' in r.text or '분석' in r.text.strip():
                    if len(r.text.strip()) < 10:
                        r.text = '실증 분석 결과'
                        r.font.name = FONT_KR

# === Slide 27: 기술통계 ===
s = prs.slides[26]
replace_content(s, '기술통계 및 지역별 분포', [
    '[전체 표본] n = 391,826건 (2019.1~2025.12, 215개 행정동)',
    '',
    '[종속변수: ㎡당 가격 (만원/㎡)]',
    '• 평균 1,337.7 / 중앙값 1,142.7 / 최대 10,586.7',
    '• 25%분위 844.4 / 75%분위 1,614.9',
    '• log 변환 후: 평균 7.0747, 표준편차 0.4879',
    '',
    '[지역별 단위가격]',
    '• 강남3구 (n=65,077): 평균 2,207.5만원/㎡',
    '• 비강남 (n=326,749): 평균 1,164.5만원/㎡',
    '• 배율 1.90배 (총가격 2.20배 → 면적효과 일부 분해)',
    '',
    '[독립변수 분포]',
    '• 전용면적 평균 75.76㎡, 중앙값 79.47㎡',
    '• 건물연령 평균 19.85년',
    '• 기준금리 분포 0.50%~3.50% (팬데믹·긴축·인하 전환 포함)',
])

# === Slide 28: 상관관계 + VIF ===
s = prs.slides[27]
remove_shapes_matching(s, ['면적당 임대료', '5개권역', '청년주택'], max_len=500)
replace_content(s, '상관관계 분석 (log ㎡당 가격 기준)', [
    '[양의 상관관계 상위]',
    '• 강남구분 r = +0.481 (가장 강한 연관)',
    '• 백화점수 r = +0.396 (상업 집적)',
    '• M2 통화량 r = +0.353 (유동성)',
    '• 소비자물가지수 r = +0.331 (인플레이션)',
    '• 층 r = +0.187',
    '',
    '[약한 상관 / 음의 상관]',
    '• 전용면적 r = +0.053 (총가격 0.581 → 규모효과 분리 증거)',
    '• 건물연령 r = −0.082 (재건축 U자형 비선형 배후)',
    '• 어린이집수 r = −0.296 (비강남 주거밀집지 혼재)',
    '',
    '[VIF Top 5]',
    '• 기준금리 110.11, CD금리 96.44, CPI 42.87, M2 28.38',
    '• 어린이집 5.99, 나머지 미시변수 모두 <5 (허용)',
])

# === Slide 29: Ablation + Moran's I ===
s = prs.slides[28]
remove_shapes_matching(s, ['권역', '연도별 면적당'], max_len=400)
replace_content(s, 'Ablation 및 Moran\'s I 강건성 점검', [
    '[Ablation: 학원수 · 어린이집수 제거]',
    '• OLS ΔR² = −0.0542 (선형 모형이 프록시에 의존)',
    '• Random Forest ΔR² = −0.0062',
    '• XGBoost ΔR² = −0.0004 (영향 거의 없음)',
    '',
    '→ XGB는 두 프록시 없이도 동등 예측 구조 유지',
    '→ 구 단위 배분 한계가 전체 결과를 왜곡하지 않음',
    '',
    '[Moran\'s I 잔차 공간 자기상관, 같은 구 내 인접 가중]',
    '• OLS: I = 0.3247, Z = 8.882, p < 0.001  (강한 양의 자기상관)',
    '• XGBoost: I = 0.0042, Z = 0.252, p = 0.880  (통계적 유의성 없음)',
    '',
    '→ XGB의 비선형 조합이 공간 구조 상당 부분 흡수',
    '→ 다만 정밀 경계 가중 행렬 재검증은 후속 과제',
])

# === Slide 30: 모형 성능 비교 (표) ===
s = prs.slides[29]
replace_content(s, '모형별 · 분할별 예측 성능 비교', [
    '[분할 | OLS | Random Forest | XGBoost] (R²_log / Median APE)',
    '',
    '무작위 : 0.5061 / 22.04%  |  0.8937 / 7.51%  |  0.9554 / 4.16%',
    'Group  : 0.4618 / 22.35%  |  0.7414 / 13.12% |  0.8005 / 11.72%',
    '시간순 : 0.3939 / 24.74%  |  0.6955 / 14.82% |  0.7972 / 12.61%',
    '',
    '[핵심 해석]',
    '• 모든 분할에서 XGB > RF >> OLS 일관',
    '• 무작위 R²=0.9554: 동일 단지 내 예측 성능',
    '• Group / 시간순 R²≈0.80: 미경험 단지·미래 시점 일반화',
    '• Group ≈ 시간순 → 단지 누수 > 국면 이동 효과',
    '',
    '[일반화 참조값]',
    '• 미경험 단지 예측: R²≈0.80, Median APE≈12% (실무 AVM 기준선)',
])

# === Slide 31: SHAP 결과 + 지역 비교 ===
s = prs.slides[30]
replace_content(s, 'SHAP 핵심 결과 및 지역 별도 비교', [
    '[전체 모형 SHAP Top 6 — 규모 정규화 후 재편]',
    '• 건물연령 12.43% (1위, 총가격에서 3위)',
    '• 강남구분 11.74%',
    '• 어린이집수 11.53% ⭐ (총가격 8위 → 3위)',
    '• M2 통화량 10.77%',
    '• 학원수 9.61% ⭐ (총가격 7위 → 5위)',
    '• 전용면적 9.55% (총가격 1위 22.2% → 6위, 규모효과 분리)',
    '',
    '[지역 별도 모형 SHAP Top 3 비교]',
    '• 강남3구 (R²=0.9356): 건물연령16.5% - CPI 13.3% - 백화점수 13.0%',
    '• 비강남 (R²=0.9421): 건물연령14.8% - CPI 11.8% - 전용면적 11.2%',
    '• 비강남 Top 4~6: 어린이집 10.57% · M2 10.11% · 학원수 8.63%',
    '',
    '→ 강남 = 재건축·상업 주도 / 비강남 = 주거권·사교육 주도',
])
# SHAP bar plot 삽입
shap_bar = os.path.join(PLOTS, 'fig5_shap_bar.png')
if os.path.exists(shap_bar):
    pic = s.shapes.add_picture(shap_bar, Inches(12), Inches(7), height=Inches(4))

# === Slide 32 (섹션 표지 05) ===
s = prs.slides[31]
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if '결론' in r.text and len(r.text.strip()) < 10:
                    r.text = '결론 및 시사점'
                    r.font.name = FONT_KR

# === Slide 33: 시사점 ===
s = prs.slides[32]
replace_content(s, '시사점', [
    '[학술적 시사점]',
    '• 국내 아파트 시장에 XAI 체계적 적용 범위 확장',
    '• 단위면적 정규화 → 가려진 질적·입지적 설명 신호 드러냄',
    '• 지역 별도 모형으로 공간 이질성 직접 실증 (통합 SHAP 재집계 대비)',
    '',
    '[정책 참고점]',
    '• 강남 vs 비강남: 단위가격 결정 구조가 질적으로 상이',
    '• 재건축 규제는 강남에, 주거·교육 인프라는 비강남에 차별적 효과',
    '• M2·CPI 거시 변수는 시장 국면 대리신호 (인과 근거 ✗)',
    '• 건물연령 U자형 패턴 → 재건축 정책 연령 구간별 차등 검토',
    '',
    '[실무적 시사점 (AVM)]',
    '• 단위가격 SHAP 분해는 면적에 가려지지 않은 질적 프리미엄 투명화',
    '• 프롭테크·감정평가 실무 XAI 설명 책무 지원',
])

# === Slide 34: 한계 및 향후 과제 ===
s = prs.slides[33]
remove_shapes_matching(s, ['공공임대주택 정책', '향후 연구에서는 이러한 한계를 보완'], max_len=400)
replace_content(s, '연구의 한계 및 향후 과제', [
    '[한계]',
    '• 시설 변수 stock 정보 병합 (연도별 변동 미반영) ← 교수 지적',
    '• 학원수·어린이집수 구 단위 균등 배분 프록시 한계',
    '• 질적 변수 (브랜드·조망·재건축 단계) 미반영',
    '• Moran\'s I 간이 가중 행렬 (정밀 경계 기반 미적용)',
    '• 공간계량 모형(SAR/SEM) 직접 비교 미수행',
    '• 강남3구 표본 규모 제약 (비강남의 1/5)',
    '',
    '[향후 과제]',
    '• 연도별 시설 패널 구축 (행안부·교육부 연간 스냅샷)',
    '• GIS 기반 최근접 거리(nearest distance) 변수',
    '• Queen/Rook 경계 기반 정밀 공간 가중 행렬 재검증',
    '• LightGBM · CatBoost · TabNet 추가 비교',
    '• LIME · ALE · Permutation Importance 교차 XAI 검증',
    '• 수도권·지방 대도시로 분석 범위 확장',
])

# === Slide 35 (참고문헌) 유지
# === Slide 36 (감사합니다)
s = prs.slides[35]
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    r.font.name = FONT_KR

# ========== 전역 폰트 정규화 (한글 깨짐 방지) ==========
print("전역 폰트 정규화 중...")
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        # 숫자만 있는 경우 Montserrat 유지, 그 외 한글/한영혼합은 Pretendard
                        has_korean = any('ㄱ' <= c <= '힣' for c in r.text)
                        if has_korean:
                            r.font.name = FONT_KR
                            rPr = r._r.get_or_add_rPr()
                            for existing in rPr.findall(qn('a:ea')):
                                rPr.remove(existing)
                            ea = etree.SubElement(rPr, qn('a:ea'))
                            ea.set('typeface', FONT_KR)

prs.save(DST)
print(f"생성 완료: {DST}")
print(f"파일 크기: {os.path.getsize(DST)/1024/1024:.1f} MB")
print(f"슬라이드 수: {len(prs.slides)}")
