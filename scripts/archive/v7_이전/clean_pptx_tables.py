#!/usr/bin/env python3
"""content 슬라이드의 양재영 원본 table/picture 잔재 제거 + 본문 영역 확장"""
import os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

DST = '/Users/seokmogu/project/thesis-seoul-apt-xai/paper/중간발표_v7_m2price.pptx'
PLOTS = '/Users/seokmogu/project/thesis-seoul-apt-xai/results/plots_v7_m2price'

prs = Presentation(DST)

# 섹션 표지 (제거하지 않음): Slide 3(index 2), 8, 18, 26, 32
SECTION_COVER_IDX = {2, 7, 17, 25, 31}
# 목차: Slide 2 (index 1)
TOC_IDX = {1}
# 표지: Slide 1 (index 0)
TITLE_IDX = {0}
# 참고문헌/감사: Slide 35, 36
KEEP_ORIGINAL_IDX = {34, 35}
# 이 중 어느 것도 아닌 슬라이드 = content slide
PROTECT_IDX = SECTION_COVER_IDX | TOC_IDX | TITLE_IDX | KEEP_ORIGINAL_IDX

def slide_area(prs):
    return (prs.slide_width or 1) * (prs.slide_height or 1)

total_area = slide_area(prs)
removed_tables = 0
removed_pics = 0

for i, slide in enumerate(prs.slides):
    if i in PROTECT_IDX:
        continue
    # content slide: 모든 table 제거 + 큰 picture 제거
    to_remove = []
    for sh in slide.shapes:
        if sh.has_table:
            to_remove.append(('table', sh))
        elif str(sh.shape_type) == 'PICTURE (13)':
            # 슬라이드의 5% 이상 차지하는 picture는 원본 표/그림으로 간주해 제거
            area = (sh.width or 0) * (sh.height or 0)
            if area > 0.02 * total_area:
                to_remove.append(('pic', sh))
    for kind, sh in to_remove:
        try:
            sh._element.getparent().remove(sh._element)
            if kind == 'table': removed_tables += 1
            else: removed_pics += 1
        except: pass

print(f"제거: tables={removed_tables}, pictures={removed_pics}")

# 본문 텍스트 영역 확장 + 재정렬: content slide에서 본문 박스가 있으면
# 크기를 슬라이드 가로 90% 세로 70%로 확장
W, H = prs.slide_width, prs.slide_height
from pptx.util import Inches as I
for i, slide in enumerate(prs.slides):
    if i in PROTECT_IDX:
        continue
    # 가장 큰 text_frame (본문으로 추정) 찾기
    candidates = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.strip()
            if t in ['1절','2절','3절','4절','5절','1','2','3','4','5','6','7','8','9','10',
                    '11','12','13','14','15','16','17','18','19','20','21','22','23','24',
                    '25','26','27','28','29','30','31','32','33','34','35','36']:
                continue
            # 본문으로 추정되는 긴 박스
            if len(t) > 50:
                candidates.append(sh)
    if not candidates:
        continue
    candidates.sort(key=lambda s: (s.width or 0)*(s.height or 0), reverse=True)
    main = candidates[0]
    # 영역 확장
    main.left = I(0.8)
    main.top = I(1.3)
    main.width = W - I(1.6)
    main.height = H - I(2.0)

# SHAP 그림은 슬라이드 31에만 추가 (이미 있으면 중복 방지)
s = prs.slides[30]
has_shap = any(str(sh.shape_type) == 'PICTURE (13)' for sh in s.shapes)
shap_bar = os.path.join(PLOTS, 'fig5_shap_bar.png')
if not has_shap and os.path.exists(shap_bar):
    s.shapes.add_picture(shap_bar, I(12), I(7), height=I(4))

prs.save(DST)
print(f"저장: {DST}")