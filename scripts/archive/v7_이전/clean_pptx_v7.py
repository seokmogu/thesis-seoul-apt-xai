#!/usr/bin/env python3
"""PPTX v7 클린업: 잔존 원본(청년안심주택) 텍스트 제거 + 표지 재정리"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

DST = '/Users/seokmogu/project/thesis-seoul-apt-xai/paper/중간발표_v7_m2price.pptx'
PLOTS = '/Users/seokmogu/project/thesis-seoul-apt-xai/results/plots_v7_m2price'
CYAN = RGBColor(0x08, 0xA5, 0xC1)
DARK = RGBColor(0x39, 0x39, 0x39)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(DST)

# 제거할 잔존 패턴
REMOVE_PATTERNS = [
    '1인 가구 연립/다세대 주택',
    '-1인 가구 연립/다세대 주택을 중심으로-',
    '조 미 정',
    '양 재 영',
    '도시부동산개발전공',
    '이 공급되는 지역의 전월세 가격 변화를 살펴봄으로써 해당 정책의 효과성 검증',
    '향후 연구에서는 이러한 한계를 보완하여, 공공임대주택 정책',
    '공간적범위: 서울특별시 5대권역',
    '청년안심주택',
    '1인 가구',
    '면적 기준',
    '<표 1> 1인 가구 비율',
    '<표 2> 거처 종류별',
    '의 전월세',
    '임대주택',
    '전월세 가격',
]
# 제거 후 새로 넣을 컨텐츠
REPLACE_MAP = [
    ('이 공급되는 지역의 전월세 가격 변화를 살펴봄으로써 해당 정책의 효과성 검증',
     '을 통해 규모효과를 정규화하고 질적·입지적 설명 신호를 드러내고자 함.'),
    ('1인 가구 연립/다세대 주택을 중심으로',
     '단위면적당 가격·행정동·지역 별도 모형'),
    ('-1인 가구 연립/다세대 주택을 중심으로-',
     '- 단위면적당 가격·행정동·지역 별도 모형 -'),
    ('조 미 정', '최 **'),
    ('양 재 영', '박 현 근'),
    ('도시부동산개발전공', '도시부동산정책전공'),
    ('공간적범위: 서울특별시 5대권역', '공간범위: 서울 25개 자치구·215개 행정동'),
    ('향후 연구에서는 이러한 한계를 보완하여, 공공임대주택 정책의 효과를 더욱 정교하고 포괄적으로 분석함으로써 학술적·정책적 기여를 확대할 필요가 있',
     '향후 연구에서는 연도별 시설 패널·GIS 거리 변수·공간계량 모형을 결합하여 단위가격 설명 구조 분석을 심화할 필요가 있'),
]

def replace_run_text(run, old, new):
    if old in run.text:
        run.text = run.text.replace(old, new)
        return True
    return False

def remove_shape_if_matches(slide, patterns):
    """shape 전체 텍스트가 제거 대상 패턴만 담고 있으면 제거"""
    to_remove = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt: continue
        # 전체가 제거 대상 패턴과 강하게 매칭되면 제거
        for pat in patterns:
            if pat in txt and len(txt) < len(pat) + 30:  # 거의 그 문자만 있으면
                to_remove.append(sh)
                break
    for sh in to_remove:
        try:
            sp = sh._element
            sp.getparent().remove(sp)
        except: pass

# 1. 표지(Slide 1) 대대적 정리
slide1 = prs.slides[0]
# 모든 shape 조사
for sh in list(slide1.shapes):
    if sh.has_text_frame:
        tf = sh.text_frame
        full = tf.text
        for para in tf.paragraphs:
            for run in para.runs:
                for old, new in REPLACE_MAP:
                    replace_run_text(run, old, new)

# 2. 전체 슬라이드에 REPLACE_MAP 적용
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in REPLACE_MAP:
                        replace_run_text(run, old, new)

# 3. 특정 슬라이드(4, 5)에서 <표 1>, <표 2> 레퍼런스 제거
for idx in [3, 4, 13, 14, 15, 16]:
    if idx < len(prs.slides):
        slide = prs.slides[idx]
        for sh in list(slide.shapes):
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if ('<표' in txt and len(txt) < 50) or ('1인 가구' in txt) or ('거처 종류' in txt):
                    try:
                        sh._element.getparent().remove(sh._element)
                    except: pass

# 4. 슬라이드 1에 깔끔한 제목 강제 설정
slide1 = prs.slides[0]
# 추가 정보 박스
for sh in list(slide1.shapes):
    if sh.has_text_frame:
        txt = sh.text_frame.text
        # "이 인근 전월세" 같은 문구를 단순 정리
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if '이 인근 전월세 가격에 미치는 영향 연구' in run.text:
                    run.text = run.text.replace('이 인근 전월세 가격에 미치는 영향 연구', '')
                if '인근 전월세 가격에 미치는 영향 연구' in run.text:
                    run.text = run.text.replace('인근 전월세 가격에 미치는 영향 연구', '')

# 5. 슬라이드 31에 SHAP 막대 그림 추가 (이미 추가된 상태일 수 있음, 중복 방지)
if len(prs.slides) >= 31:
    slide = prs.slides[30]
    has_pic = any(sh.shape_type == 13 for sh in slide.shapes)
    shap_bar = os.path.join(PLOTS, 'fig5_shap_bar.png')
    if not has_pic and os.path.exists(shap_bar):
        from pptx.util import Emu
        slide.shapes.add_picture(shap_bar, Emu(10*914400), Emu(2*914400), height=Emu(7*914400))

prs.save(DST)
print(f"클린업 완료: {DST}")

# 검증
prs2 = Presentation(DST)
for i in [0, 3, 4, 33]:
    s = prs2.slides[i]
    print(f"\n--- Slide {i+1} ---")
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.replace('\n', ' | ')[:150]
            print(f"  {t}")
