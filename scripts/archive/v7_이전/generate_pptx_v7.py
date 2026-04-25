#!/usr/bin/env python3
"""논문 v7(단위면적당 가격) 중간발표용 PPTX 생성.
양재영_중간발표_1203.pptx 템플릿의 레이아웃·색상·폰트를 재활용하되,
본 논문 내용으로 슬라이드별 텍스트를 재작성한다."""
import os, sys, shutil, copy
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = '/Users/seokmogu/Downloads/양재영_중간발표_1203.pptx'
DST = '/Users/seokmogu/project/thesis-seoul-apt-xai/paper/중간발표_v7_m2price.pptx'
PLOTS = '/Users/seokmogu/project/thesis-seoul-apt-xai/results/plots_v7_m2price'

shutil.copy(SRC, DST)
prs = Presentation(DST)

CYAN = RGBColor(0x08, 0xA5, 0xC1)
DARK = RGBColor(0x39, 0x39, 0x39)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 슬라이드별 원본 텍스트 → 새 텍스트 매핑 (논문 v7 기준)
# 각 슬라이드의 기존 텍스트를 부분 매칭으로 찾아서 replace
REPLACEMENTS = {
    # Slide 1 (표지)
    '서울시 청년안심주택이 인근 전월세 가격에 미치는 영향 연구':
        'XGBoost와 SHAP을 활용한 서울시 아파트 단위면적당 매매가격의 설명 패턴 분석',
    '논문발표': '석사학위 중간발표',
    '양재영': '박현근',
    # Slide 2 목차
    '서울시 1인 가구 증가': '부동산=한국 가계자산 75.2%: 서울 아파트가 자산·주거 안정의 핵심',
    '청년안심주택의 준공이 인근 비(非) 아파트 유형의 주택 전월세에 영향을 미치는지 분석하고자 함':
        '단위면적당 거래가격(log ㎡당 만원) 예측과 SHAP 기반 설명 구조 분석: 규모효과 정규화로 질적 신호 부각',
    '공간적범위: 서울특별시 5대권역 내':
        '공간적 범위: 서울 215개 행정동 | 시간적 범위: 2019.1~2025.12 (총 391,826건)',
    '청년안심주택': '헤도닉 가격모형 / XGBoost / SHAP / XAI',
    '1인 가구 면적 기준': '18개 변수 구성 (물리·입지·환경·거시)',
    '선행연구고찰': '선행연구 및 차별성',
    '연구의 흐름': '연구 흐름: 데이터 → 전처리 → 모형 → SHAP → 지역별 → 시사점',
    '연구방법 및 변수설정': '변수 설정 및 하이퍼파라미터',
    '분석방법론': '분석 방법: OLS + Random Forest + XGBoost + SHAP',
    '기술통계량': '기술통계 (단위면적당 가격 기준)',
    '상관관계분석': '상관관계 분석',
    '차이분석(ANOVA 분석)': 'VIF 진단 및 Ablation',
    '다중회귀분석': '모형 성능 비교 및 SHAP 중요도',
    '분석결과': '지역 별도 모형 및 시기별 분석 결과',
    '시사점': '학술·정책·실무 시사점',
    '한계점 및 향후 연구방향': '연구의 한계 및 향후 과제',
    '참고문헌': '참고문헌',
}

# 각 슬라이드 본문을 구체적으로 재구성
# (슬라이드 번호: (원본 텍스트 패턴, 새 텍스트) 리스트)
SLIDE_CONTENTS = {
    # Slide 4: 연구배경
    4: {
        'full_replace': {
            # 기존 긴 설명을 새 배경으로
            '서울시 1인 가구 증가': '연구의 배경',
        },
        'bullets': [
            '• 한국 가계 자산의 75.2%가 실물자산, 대부분 부동산 (통계청, 2024)',
            '• 서울 아파트: 자산·주거 안정의 핵심 + 거시경제 파급효과',
            '• 기존 OLS 헤도닉 모형은 비선형·상호작용 포착 한계',
            '• ML(RF·XGBoost)은 성능 우수하나 블랙박스 문제',
            '• XAI(SHAP)로 해석 가능성 확보 → 국내 아파트 시장 체계적 적용 부족',
        ]
    },
    # Slide 5: 연구 필요성
    5: {
        'full_replace': {
            '청년안심주택의 준공이 인근 비(非) 아파트 유형의 주택 전월세에 영향을 미치는지 분석하고자 함': '연구의 필요성 및 목적',
        },
        'bullets': [
            '[연구 동기] 왜 이 주제인가?',
            '• AVM 실무 확산 vs 국내 설명 프레임워크 부재',
            '• 기존 한국 ML 연구의 3공백',
            '  (a) 자치구 단위 집계 (b) 단일 모형 비교 (c) 통합 SHAP만',
            '• 총가격 모형은 면적이 SHAP 상위 독점 → 질적 신호 가려짐',
            '',
            '[본 연구 차별성]',
            '① 단위면적당 가격 정규화 (규모효과 분리)',
            '② 행정동(215개) 단위 세분화',
            '③ 지역 별도 모형 (전체·강남·비강남) 직접 비교',
        ]
    },
    # Slide 6: 연구 범위
    6: {
        'full_replace': {
            '공간적범위: 서울특별시 5대권역 내': '연구의 범위',
        },
        'bullets': [
            '• 공간: 서울 25개 자치구 = 215개 행정동',
            '• 시간: 2019.1 ~ 2025.12 (84개월)',
            '• 대상: 아파트 매매 실거래 391,826건',
            '• 데이터 출처: 국토교통부 실거래가 / 서울열린데이터광장 / NEIS / 한국은행 ECOS',
            '• 종속변수: log(거래금액 / 전용면적) — 단위면적당 가격',
            '• 독립변수: 18개 (물리 3 + 입지 2 + 환경 9 + 거시 4)',
        ]
    },
    # Slide 9-12: 개념 설명 (헤도닉·ML·XAI)
    9: {
        'bullets': [
            '헤도닉 가격모형 (Hedonic Price Model)',
            '• Lancaster(1966), Rosen(1974)의 이론적 기초',
            '• 이질적 재화의 가격 = f(속성₁, 속성₂, ..., 속성ₙ)',
            '• 전통적 추정: 다중회귀분석(OLS)',
            '• 한계: 선형 가정, 상호작용 미반영, 다중공선성 취약',
            '',
            '단위면적당 가격 정규화 (본 연구 설계)',
            '• Y = log(P/A) = ln(거래금액/전용면적)',
            '• Malpezzi(2003), Sirmans et al.(2005) 표준 log-price 전통 연결',
            '• 준 탄력성 해석: β → (e^β - 1)·100% 변화',
        ]
    },
    10: {
        'bullets': [
            '머신러닝 기반 예측',
            '• Random Forest (Breiman, 2001): 배깅 기반 독립 트리 앙상블',
            '• XGBoost (Chen & Guestrin, 2016): 그래디언트 부스팅 + 정규화',
            '• 비선형 관계·고차 상호작용 자동 학습',
            '',
            '하이퍼파라미터 (본 연구)',
            '• XGBoost: max_depth=8, learning_rate=0.1, min_child_weight=5,',
            '  reg_alpha=0.1, reg_lambda=1.0, tree_method=hist',
            '• Random Forest: n_estimators=200, max_depth=15, min_samples_leaf=5',
            '• 공통: 무작위 70/10/20 분할, early stopping(50) 적용',
        ]
    },
    11: {
        'bullets': [
            'XAI (eXplainable AI) & SHAP',
            '• 블랙박스 모형의 해석 가능성 확보 수단',
            '• SHAP (Lundberg & Lee, 2017)',
            '  - 게임이론 Shapley value를 예측값 분해에 적용',
            '  - 글로벌 중요도 + 개별 예측 해석 + 의존성(Dependence) 분석',
            '• TreeSHAP (Lundberg et al., 2020)',
            '  - 트리 기반 모형 전용 정확 계산 알고리즘',
            '',
            '본 연구 활용',
            '• SHAP 변수 중요도 (글로벌)',
            '• SHAP Dependence Plot (비선형 검증)',
            '• SHAP Force Plot (개별 거래 분해)',
        ]
    },
    12: {
        'bullets': [
            '국내 선행연구 한계',
            '• 자치구(25개) 단위 집계 → 세분화 부족',
            '• 단일 모형 예측 성능 비교에 머무름',
            '• 통합 모형 SHAP 중요도만 보고 → 지역 이질성 미실증',
            '',
            '해외 선행연구 (XAI 적용)',
            '• Mora-García et al.(2022), Neves et al.(2024): 유럽 사례',
            '• Choy & Ho(2023): 체계적 문헌 리뷰',
            '• Chun et al.(2025), Kim et al.(2025), An et al.(2025): 한국 초기 적용',
            '',
            '본 연구 차별성',
            '• 단위면적 정규화 + 행정동 + 지역 별도 모형 (3축 동시)',
        ]
    },
    # Slide 14-17: 선행연구
    14: {'bullets': [
        '선행연구 1: 머신러닝 예측 성능 비교',
        '• 김이환 외(2022): 기계학습 방법론 vs 헤도닉, 우수성 확인',
        '• 김학현 외(2023): DNN·XGBoost·CatBoost 비교',
        '• 이선구(2025): XGBoost AVM 은평구 다세대 86.4%',
        '• 배성완·유정석(2018): 초기 ML 예측 연구',
        '• Čeh et al.(2018): RF vs OLS 우위 실증',
    ]},
    15: {'bullets': [
        '선행연구 2: XAI·SHAP 적용',
        '• Neves et al.(2024): 리스본 스마트시티, 공원거리 SHAP',
        '• Mora-García et al.(2022): COVID-19 시기 ML+DL 비교',
        '• Chun et al.(2025): 서울 XAI 최초 체계적 적용',
        '• Kim et al.(2025): 설명가능 AVM 한국 주택시장',
        '• An et al.(2025): 헤도닉 vs ML 투명성 비교',
    ]},
    16: {'bullets': [
        '선행연구 3: 공간·거시 통합',
        '• Kim et al.(2022): 공간회귀분석 공간의존성',
        '• 배종찬·정재호(2021): VEC 거시+부동산정책',
        '• 노산하·신진호(2021): 검색빈도 주택가격',
        '• 정진오·정재호(2023): 조세·금융정책 차별효과',
        '• 김명진·서원석(2023): 코로나19 소매시설 시계열',
    ]},
    17: {'bullets': [
        '본 연구의 차별성',
        '',
        '[방법론]',
        '• 단위면적당 가격 (log) 정규화 → 규모효과 분리',
        '• 행정동 215개 분석틀 (자치구 25개 대비 세분화)',
        '• 지역 별도 모형 (전체/강남/비강남 3개 독립 적합)',
        '',
        '[실증 설계]',
        '• 3가지 분할 (random/group/chronological) 동시 보고',
        '• Ablation + Moran\'s I 강건성 점검',
        '• SHAP 글로벌 + Dependence + 지역 비교',
    ]},
    # Slide 19-25: 연구방법
    19: {'bullets': [
        '연구 흐름',
        '',
        '① 데이터 수집 (국토부·서울데이터광장·NEIS·ECOS)',
        '② 법정동→행정동 매핑 (Nominatim + GeoJSON)',
        '③ 변수 구축 (18개) + log(㎡당 가격) 파생',
        '④ OLS / RF / XGBoost 적합',
        '⑤ SHAP 분석 (글로벌·Dependence·지역별)',
        '⑥ 강건성 점검 (Ablation, Moran\'s I)',
        '⑦ 지역 별도 모형 (전체/강남/비강남)',
        '⑧ 시기별 3구간 보조 분석',
    ]},
    20: {'bullets': [
        '변수 설정 (18개)',
        '',
        '[물리적 특성] 전용면적(㎡), 층, 건물연령',
        '[입지 특성] 강남구분(더미), 지하철역수',
        '[환경 특성] 초·중·고등학교수, CCTV, 백화점수,',
        '           공원수, 도서관수, 학원수*, 어린이집수*',
        '[거시경제] 기준금리, CD금리, 소비자물가지수, M2',
        '',
        '* 구 단위 집계를 행정동 수로 균등 배분한 프록시',
        '* 시설 변수: 수집시점 stock 정보 (연도별 변동 미반영 — 한계)',
    ]},
    21: {'bullets': [
        '분할 방식 3종',
        '',
        '① 무작위 분할 (70/10/20)',
        '   • 동일 단지 내 예측 성능 지표',
        '   • 반복거래 누수 가능성',
        '',
        '② Group 분할 (법정동+아파트명 기준, 8,601개 단지)',
        '   • 미경험 단지 일반화 성능',
        '   • 단지 반복거래 누수 차단',
        '',
        '③ 시간순 분할 (학습<2024.7, 테스트 2024.7~2025.12)',
        '   • 미래 시점 예측 성능',
        '   • 시장 국면 이동 반영',
    ]},
    22: {'bullets': [
        '모형 구성',
        '',
        '• OLS (다중회귀): 베이스라인, 선형 가정',
        '  y = β₀ + Σβᵢxᵢ + ε   (y = log(P/A))',
        '',
        '• Random Forest: 200 trees, depth 15',
        '  배깅 기반 독립 트리 앙상블',
        '',
        '• XGBoost: 2000 rounds, depth 8, lr 0.1',
        '  그래디언트 부스팅 + L1/L2 정규화',
        '  early stopping patience=50',
        '',
        '• 모든 실험 동일 하이퍼파라미터 유지',
    ]},
    23: {'bullets': [
        '평가 지표',
        '',
        '• R²_log: log 스케일 결정계수 (모형 간 동일 척도)',
        '• RMSE_log: log 스케일 제곱근 오차',
        '• MAE_log: log 스케일 절대 오차',
        '• MAPE: 원 스케일(만원/㎡) 평균 절대 상대 오차',
        '• Median APE: 원 스케일 중앙값 절대 상대 오차',
        '',
        '[강건성 점검]',
        '• Ablation (학원수·어린이집수 제거 영향)',
        '• Moran\'s I (잔차 공간 자기상관)',
        '• 월 클러스터-강건 표준오차',
    ]},
    24: {'bullets': [
        'SHAP 분석 방법',
        '',
        '• TreeSHAP (Lundberg et al., 2020)',
        '  - XGBoost 전용 정확 Shapley value',
        '',
        '• 샘플 5,000건 (테스트 데이터에서 무작위)',
        '',
        '• 글로벌 중요도: mean(|SHAP_log|) 기준 순위',
        '• Dependence Plot: 변수값 vs SHAP값 비선형 패턴',
        '• Force Plot: 개별 거래 기여도 분해',
        '',
        '• 역변환 해석: log SHAP → exp(SHAP) ≈ ±x% 변화',
    ]},
    25: {'bullets': [
        '지역 별도 모형 설계 (교수 피드백 반영)',
        '',
        '• 기존 접근: 통합 모형 + SHAP 지역별 재집계',
        '• 본 연구: 3지역 독립 XGBoost 적합',
        '',
        '① 전체 (17 features, 강남더미 제거): n=391,826',
        '② 강남3구: n=65,077 (강남·서초·송파)',
        '③ 비강남: n=326,749',
        '',
        '• 동일 하이퍼파라미터, 동일 70/10/20 분할',
        '• 지역 내부 SHAP 구조 직접 비교',
        '',
        '→ 연도별 분해 대신 공간 이질성 실증에 초점',
    ]},
    # Slide 27-31: 분석결과
    27: {'bullets': [
        '기술통계량 요약',
        '',
        '• 총 391,826건 (215개 행정동, 2019.1~2025.12)',
        '• 거래금액: 평균 10.3억, 중앙 8.3억, 최대 290억',
        '• ㎡당 가격: 평균 1,338만원/㎡, 중앙 1,143만원/㎡',
        '• 전용면적: 평균 75.8㎡, 중앙 79.5㎡',
        '• 건물연령: 평균 19.9년',
        '',
        '[지역별]',
        '• 강남3구: 평균 2,208만원/㎡ (n=65,077, 17%)',
        '• 비강남: 평균 1,165만원/㎡ (n=326,749, 83%)',
        '• 배율: 1.90배 (총가격 2.20배보다 작음 → 면적효과 일부 분해)',
    ]},
    28: {'bullets': [
        '상관관계 분석 (log ㎡당 가격 기준)',
        '',
        '• 강남구분 r = +0.481 ← 가장 강한 양의 상관',
        '• 백화점수 r = +0.396',
        '• M2 r = +0.353',
        '• 소비자물가지수 r = +0.331',
        '• 층 r = +0.187',
        '• 전용면적 r = +0.053 ← 총가격 0.581 → 정규화 후 급감',
        '• 건물연령 r = -0.082',
        '• 어린이집수 r = -0.296 (비강남 밀집 혼재)',
        '',
        '[총가격 vs 단위가격 차이]',
        '면적-가격 상관 0.581 → 0.053: 규모효과 분리의 정량적 증거',
    ]},
    29: {'bullets': [
        'VIF 진단 및 Ablation',
        '',
        '[VIF Top 5]',
        '• 기준금리 110.11, CD금리 96.44, 소비자물가 42.87, M2 28.38',
        '• 어린이집수 5.99, 중학교 4.13, 초등학교 4.04',
        '• 나머지 미시 변수는 모두 < 5 (허용)',
        '',
        '[Ablation: 학원수·어린이집수 제거]',
        '• OLS ΔR² = -0.0542 (선형 모형이 프록시에 의존)',
        '• RF ΔR² = -0.0062',
        '• XGB ΔR² = -0.0004 (영향 거의 없음)',
        '',
        '→ XGB는 프록시 없이도 동등 예측 구조 유지',
    ]},
    30: {'bullets': [
        '모형 성능 비교 (세 가지 분할)',
        '',
        '[무작위 분할] XGB 우위',
        '• OLS R²_log=0.5061 | RF 0.8937 | XGB 0.9554 (MedAPE 4.16%)',
        '',
        '[Group 분할] 단지 반복거래 누수 차단',
        '• OLS 0.4618 | RF 0.7414 | XGB 0.8005 (MedAPE 11.72%)',
        '',
        '[시간순 분할] 2024H2+ 일반화',
        '• OLS 0.3939 | RF 0.6955 | XGB 0.7972 (MedAPE 12.61%)',
        '',
        '• 모든 분할에서 XGB > RF >> OLS 일관',
        '• Group ≈ 시간순 → 국면이동보다 단지누수가 주 영향',
    ]},
    31: {'bullets': [
        'SHAP 글로벌 중요도 및 지역 별도 비교',
        '',
        '[전체 Top 5 (총가격 → 단위가격)]',
        '• 전용면적 22.2% → 건물연령 12.4%',
        '• 강남구분 16.0% → 강남구분 11.7%',
        '• 건물연령 11.4% → 어린이집수 11.5% ⭐',
        '• M2 8.3% → M2 10.8%',
        '• 백화점수 6.8% → 학원수 9.6% ⭐',
        '',
        '[지역 별도 모형 SHAP Top 3]',
        '• 강남: 건물연령(16.5%) - 소비자물가(13.3%) - 백화점(13.0%)',
        '• 비강남: 건물연령(14.8%) - 소비자물가(11.8%) - 전용면적(11.2%)',
        '→ 강남=재건축/상업, 비강남=주거권/사교육 주도',
    ]},
    # Slide 33-34: 시사점·한계
    33: {'bullets': [
        '시사점',
        '',
        '[학술적]',
        '• 국내 아파트 XAI 실증적 적용 범위 확장',
        '• 단위면적 정규화로 숨은 질적 설명 신호 드러냄',
        '• 지역 별도 모형으로 공간 이질성 직접 실증',
        '',
        '[정책적]',
        '• 강남/비강남은 단위가격 결정 구조가 상이',
        '• 재건축 규제는 강남에, 주거권 인프라는 비강남에 차별적 효과',
        '• M2·CPI 거시변수는 시장 국면 대리신호로 활용 (인과 근거 아님)',
        '',
        '[실무적]',
        '• AVM 설명 책무: 단위가격 SHAP 분해가 감정평가 이해도 제고',
        '• 프롭테크 XAI 기능 구현 참고',
    ]},
    34: {'bullets': [
        '연구의 한계 및 향후 과제',
        '',
        '[한계]',
        '• 시설 변수 stock 정보 병합 (연도별 변동 미반영) ← 교수 지적',
        '• 학원·어린이집 구 단위 균등 배분 프록시 한계',
        '• 질적 변수 (브랜드·조망·재건축 단계) 미반영',
        '• 공간계량 모형(SAR/SEM) 직접 비교 미수행',
        '',
        '[향후 과제]',
        '• 연도별 시설 패널 구축 (행정안전부·교육부)',
        '• GIS 기반 최근접 거리(nearest distance) 변수',
        '• LightGBM·CatBoost 추가 비교',
        '• LIME·ALE 교차 XAI 검증',
        '• 수도권·지방 대도시로 범위 확장',
    ]},
    # Slide 35 참고문헌은 원본 유지
    # Slide 36 감사합니다 유지
}

def replace_text_in_shape(shape, replacements):
    """shape 안의 텍스트 교체 (부분 매칭)"""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    replaced = False
    for para in tf.paragraphs:
        for run in para.runs:
            original = run.text
            for old, new in replacements.items():
                if old in original:
                    run.text = original.replace(old, new)
                    replaced = True
                    original = run.text
    return replaced

def set_slide_bullets(slide, title_text, bullets, font_size=16):
    """슬라이드를 단순 텍스트 구조로 재작성 (title + bullets).
    기존 텍스트 박스 중 가장 큰 것을 본문으로 사용."""
    # 기존 텍스트 박스에 bullets 채우기 (기존 포맷 유지)
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    if not text_shapes:
        return
    # 가장 큰 텍스트 박스 찾기
    text_shapes.sort(key=lambda s: (s.width or 0)*(s.height or 0), reverse=True)
    main = text_shapes[0]
    # 기존 첫 줄 폰트 정보 저장
    tf = main.text_frame
    first_p = tf.paragraphs[0]
    orig_font_name = None
    orig_color = DARK
    if first_p.runs:
        r = first_p.runs[0]
        orig_font_name = r.font.name
        try:
            if r.font.color and r.font.color.rgb:
                orig_color = r.font.color.rgb
        except: pass
    # 텍스트 모두 지우기 + 새로 쓰기
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if orig_font_name:
            r.font.name = orig_font_name
        r.font.size = Pt(font_size)
        if line and not line.startswith('•') and not line.startswith(' ') and i == 0:
            r.font.bold = True
            r.font.size = Pt(font_size + 4)
            try: r.font.color.rgb = CYAN
            except: pass
        else:
            try: r.font.color.rgb = DARK
            except: pass

# 1. 공통 replacements (모든 슬라이드에서 치환)
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        replace_text_in_shape(sh, REPLACEMENTS)

# 2. 슬라이드별 세부 콘텐츠 작성
for slide_num, spec in SLIDE_CONTENTS.items():
    idx = slide_num - 1  # 0-based
    if idx >= len(prs.slides):
        continue
    slide = prs.slides[idx]
    bullets = spec.get('bullets')
    if bullets:
        # 해당 슬라이드의 가장 큰 텍스트 박스를 본문으로 재작성
        # 단, 제목 박스(상단 작은 것)는 놔두고 본문만 교체
        text_shapes = []
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip() and len(s.text_frame.text) > 10:
                text_shapes.append(s)
        if text_shapes:
            # 가장 큰 것
            text_shapes.sort(key=lambda s: (s.width or 0)*(s.height or 0), reverse=True)
            main = text_shapes[0]
            tf = main.text_frame
            # 기존 폰트 정보 보존
            first_p = tf.paragraphs[0]
            orig_font_name = None
            if first_p.runs:
                orig_font_name = first_p.runs[0].font.name
            # 텍스트 지우기
            # python-pptx는 clear()가 있음
            for p in list(tf.paragraphs):
                p._element.getparent().remove(p._element) if p._element.getparent() is not None else None
            # 재추가
            from pptx.oxml.ns import qn
            # clear 후 최소 1개 paragraph 필요
            from lxml import etree
            p_xml = etree.SubElement(tf._txBody, qn('a:p'))
            for i, line in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                r = p.add_run()
                r.text = line
                if orig_font_name:
                    r.font.name = orig_font_name
                r.font.size = Pt(14)
                if i == 0 and not line.startswith('•'):
                    r.font.bold = True
                    r.font.size = Pt(20)
                    try: r.font.color.rgb = CYAN
                    except: pass
                else:
                    try: r.font.color.rgb = DARK
                    except: pass

# 3. SHAP 그림을 슬라이드 31에 추가 (있으면)
shap_bar = os.path.join(PLOTS, 'fig5_shap_bar.png')
if os.path.exists(shap_bar) and len(prs.slides) >= 31:
    slide = prs.slides[30]  # slide 31
    pic = slide.shapes.add_picture(shap_bar, Emu(9*914400), Emu(3*914400), height=Emu(7*914400))

# 4. 저장
prs.save(DST)
print(f"생성 완료: {DST}")
print(f"파일 크기: {os.path.getsize(DST)/1024/1024:.1f} MB")
print(f"슬라이드 수: {len(prs.slides)}")
