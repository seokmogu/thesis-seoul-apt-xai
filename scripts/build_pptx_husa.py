#!/usr/bin/env python3
"""
중간발표 PPT — HUSA(한양대 도시대학원) 학술 표준 디자인 적용.
참고: HUSA 발제자료 샘플 (신형섭·박준영·김소희·한효진).

20 슬라이드:
  01 표지 / 02 CONTENTS
  03~04 INTRODUCTION (연구 배경 / 연구 질문·기여)
  05~06 LITERATURE REVIEW (선행연구·공백 / 차별성)
  07~11 METHODS (데이터·지오코딩·거리변수·시점정합·모형설계)
  12~18 RESULTS (Ablation·성능·SHAP Top·비선형·권역·연도×권역·Top1)
  19~20 CONCLUSION (기여 / 한계·마무리)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, 'figures')
OUT = os.path.join(ROOT, 'paper', '중간발표_박현근.pptx')

# === 컬러 (HUSA 학술 표준) ===
NAVY = RGBColor(0x2D, 0x3E, 0x5C)        # 메인 네이비
NAVY_LT = RGBColor(0x4B, 0x5D, 0x82)      # 라이트 네이비 (헤더 배경 등)
ACCENT = RGBColor(0xC8, 0x44, 0x3D)       # 빨강 강조
DARK = RGBColor(0x33, 0x33, 0x33)         # 본문 검정
GRAY = RGBColor(0x77, 0x77, 0x77)         # 영문 부제·페이지번호
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)        # 라인
BG_LIGHT = RGBColor(0xF1, 0xF3, 0xF6)     # 본문 박스 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

KR_FONT = 'Apple SD Gothic Neo'

EN_TITLE = 'Explanatory Patterns of Apartment Unit-Area Sale Prices in Seoul Using XGBoost and SHAP'


def set_font(run, size=14, bold=False, color=DARK, font=KR_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for ex in rPr.findall(qn('a:ea')):
        rPr.remove(ex)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    for ex in rPr.findall(qn('a:latin')):
        rPr.remove(ex)
    la = etree.SubElement(rPr, qn('a:latin'))
    la.set('typeface', font)


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, lines=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    if lines is None:
        lines = [text]
    for i, l_ in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = l_
        set_font(run, size, bold, color)
    return tb


def add_bullets(slide, l, t, w, h, bullets, size=14, color=DARK, line_spacing=1.35):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = '· ' + b
        set_font(run, size, False, color)
    return tb


def add_rect(slide, l, t, w, h, fill=NAVY, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def add_image(slide, path, l, t, w=None, h=None):
    if not os.path.exists(path):
        return None
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    return slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)


def add_chapter_header(slide, num, en_chapter, kr_subtitle, page_num):
    """HUSA 표준 챕터 헤더 — 좌상단 네이비 박스 + 한글 부제 + 우상단 영문 부제."""
    # 좌상단 네이비 박스 + 챕터 라벨 (LITERATURE REVIEW가 한 줄에 들어가도록 폭 확장)
    add_rect(slide, 0, 0, 3.2, 0.55, fill=NAVY)
    add_text(slide, 0.15, 0, 3.0, 0.55, f'{num:02d}. {en_chapter}',
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 한글 부제 (좌측)
    add_text(slide, 0.25, 0.6, 5.0, 0.4, kr_subtitle,
             size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    # 영문 부제 (우상단)
    add_text(slide, 6.4, 0.1, 6.8, 0.4, EN_TITLE,
             size=10, color=GRAY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # 헤더 하단 라인
    add_rect(slide, 0.25, 1.05, 12.85, 0.02, fill=LIGHT)
    # 우하단 페이지 번호 박스
    add_rect(slide, 12.83, 7.05, 0.5, 0.45, fill=NAVY)
    add_text(slide, 12.83, 7.05, 0.5, 0.45, f'{page_num:02d}',
             size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_notes(slide, script):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = script
    for p in tf.paragraphs:
        for r in p.runs:
            try:
                set_font(r, size=11, color=DARK)
            except Exception:
                pass


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ===== 01. 표지 =====
    s = prs.slides.add_slide(blank)
    # 좌상단 라벨
    add_rect(s, 0.4, 0.35, 2.3, 0.45, fill=NAVY)
    add_text(s, 0.4, 0.35, 2.3, 0.45, '석사학위논문 중간발표',
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 메인 네이비 박스
    add_rect(s, 0.4, 1.7, 12.5, 2.4, fill=NAVY)
    add_text(s, 0.6, 1.85, 12.1, 0.9, 'XGBoost와 SHAP을 활용한',
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 2.65, 12.1, 0.9, '서울시 아파트 단위면적당 매매가격의 설명 패턴 분석',
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 3.55, 12.1, 0.5, EN_TITLE,
             size=13, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 하단 정보
    add_text(s, 0.4, 4.5, 12.5, 0.5, '지도교수 : 고 준 호',
             size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 6.2, 12.5, 0.45, '2026. 04.',
             size=14, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 6.6, 12.5, 0.45, '한양대학교 부동산융합대학원',
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 6.95, 12.5, 0.45, '도시부동산정책전공  박  현  근',
             size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_notes(s, '''[표지]
안녕하십니까. 한양대학교 부동산융합대학원 도시부동산정책전공 박현근입니다.
"XGBoost와 SHAP을 활용한 서울시 아파트 단위면적당 매매가격의 설명 패턴 분석"이라는 주제로 석사학위논문 중간발표를 시작하겠습니다.''')

    # ===== 02. CONTENTS =====
    s = prs.slides.add_slide(blank)
    # 좌측 네이비 패널
    add_rect(s, 0, 0, 4.5, 7.5, fill=NAVY)
    add_text(s, 0.5, 1.5, 4.0, 1.0, 'CONTENTS',
             size=44, bold=True, color=WHITE, anchor=MSO_ANCHOR.TOP)
    # 우측 챕터 목록
    sections = [
        ('01', 'INTRODUCTION', '연구 배경 및 연구 질문'),
        ('02', 'LITERATURE REVIEW', '선행연구 검토와 본 연구의 차별성'),
        ('03', 'METHODS', '데이터·변수 설계·모형 구축'),
        ('04', 'RESULTS', 'Ablation·SHAP·시공간 이질성'),
        ('05', 'CONCLUSION', '결과 요약·기여·한계'),
    ]
    for i, (no, en, kr) in enumerate(sections):
        y = 1.5 + i * 0.95
        add_rect(s, 5.0, y, 0.75, 0.6, fill=NAVY_LT)
        add_text(s, 5.0, y, 0.75, 0.6, no,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 6.0, y, 7.0, 0.35, en,
                 size=18, bold=True, color=NAVY)
        add_text(s, 6.0, y + 0.35, 7.0, 0.3, kr,
                 size=12, color=GRAY)
    add_rect(s, 12.83, 7.05, 0.5, 0.45, fill=NAVY)
    add_text(s, 12.83, 7.05, 0.5, 0.45, '02', size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(s, '발표는 서론·선행연구·연구방법·결과·결론 다섯 부분으로 진행하겠습니다.')

    # ===== 03. INTRODUCTION — 연구 배경 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 1, 'INTRODUCTION', '연구 배경 : 세 가지 동기 층위', 3)
    # 3개 박스 (실무·학술·국면)
    titles = [('실무·제도', 'AVM·감정평가의\n설명 가능성 요구',
              ['자동감정평가(AVM) 실무 확산',
               '감정평가·은행·세무에서 "왜 이 가격인가" 설명 필요',
               '2022년 금리 급등 국면 모형 안정성 시험']),
              ('학술 공백', 'ML 부동산 연구의\n세 공백',
              ['공간 단위 거침(자치구·MAUP)',
               '시점 정합 부재(temporal leakage)',
               '단일 SHAP에 머무는 해석']),
              ('시장 국면', '2019~2025\n급등→조정→회복',
              ['기준금리 0.5% → 3.5% → 2.5%',
               '2022 거래량 −70% (조정기)',
               '국면 전환기 가격 신호 재구성'])]
    for i, (cat, title, items) in enumerate(titles):
        x = 0.5 + i * 4.25
        add_rect(s, x, 1.4, 4.0, 0.6, fill=NAVY_LT)
        add_text(s, x, 1.4, 4.0, 0.6, cat,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, 2.0, 4.0, 4.5, fill=BG_LIGHT)
        add_text(s, x + 0.2, 2.15, 3.6, 1.0, title,
                 size=15, bold=True, color=NAVY,
                 lines=title.split('\n'))
        add_bullets(s, x + 0.2, 3.4, 3.6, 3.0, items, size=11)
    add_text(s, 0.5, 6.65, 12.3, 0.45,
             '→ 네 가지 연구 질문 (Q1 공간 정합 · Q2 시간 정합 · Q3 시공간 이질성 · Q4 해석 프레임워크)',
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, '''[INTRODUCTION 연구 배경]
본 연구의 동기는 세 층위입니다.
실무 층위에서 AVM 확산에 따라 가격 설명 프레임워크가 요구되고 있고,
학술 층위에서는 기존 ML 부동산 연구가 공간 단위 거침, 시점 정합 부재, 단일 SHAP이라는 세 공백을 공유합니다.
시장 국면 층위에서는 2019부터 2025년까지가 급등·조정·회복을 모두 포함한 7년이라는 점입니다.
이로부터 네 가지 연구 질문을 도출했습니다.''')

    # ===== 04. INTRODUCTION — Q1~Q4 + 기여 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 1, 'INTRODUCTION', '연구 질문과 기여 3축', 4)
    # 좌측: 연구 질문 4개
    add_rect(s, 0.4, 1.3, 6.2, 0.5, fill=NAVY)
    add_text(s, 0.4, 1.3, 6.2, 0.5, '연구 질문 (Research Questions)',
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    qs = [
        ('Q1', '거리 기반 접근성 변수가 행정동 집계 대비\n예측 성능과 해석 안정성을 개선하는가?'),
        ('Q2', '연도별 시점 정합 변수가 시간역전 누수를 제거하면서\n시간순 분할에서의 일반화 성능을 유지하는가?'),
        ('Q3', '강남3구와 비강남의 가격 예측 신호는\n구성과 연도별 변화 양상이 어떻게 다른가?'),
        ('Q4', 'SHAP 권역별·연도별·면적대별 분해로 도출되는\n해석 구조가 AVM 실무에 활용 가능한가?'),
    ]
    for i, (q, body) in enumerate(qs):
        y = 1.95 + i * 1.05
        add_rect(s, 0.4, y, 0.55, 0.95, fill=NAVY_LT)
        add_text(s, 0.4, y, 0.55, 0.95, q,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.05, y + 0.05, 5.55, 0.85, body,
                 size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE,
                 lines=body.split('\n'))
    # 우측: 기여 3축
    add_rect(s, 6.9, 1.3, 6.0, 0.5, fill=NAVY)
    add_text(s, 6.9, 1.3, 6.0, 0.5, '기여 (Contributions)',
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cs = [
        ('방법론적 타당성', '행정동 집계·거리 기반·시점 정합·통합 4 시나리오 Ablation'),
        ('예측 성능 개선', '시간순 분할 R² +7.1%p (행정동 집계 대비) · 통합 모형 D 최고'),
        ('시공간 이질성 해석', '권역×연도 21 SHAP — 강남 6회 Top1 교체 vs 비강남 7년 안정'),
    ]
    for i, (t, b) in enumerate(cs):
        y = 1.95 + i * 1.4
        add_rect(s, 6.9, y, 6.0, 0.45, fill=NAVY_LT)
        add_text(s, 7.05, y, 5.85, 0.45, t,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 6.9, y + 0.55, 6.0, 0.7, b,
                 size=11, color=DARK)
    add_notes(s, '연구 질문은 공간 정합·시간 정합·시공간 이질성·해석 프레임워크 네 축이며, 기여는 방법론적 타당성·예측 성능·시공간 이질성 세 축입니다.')

    # ===== 05. LITERATURE REVIEW — 세 공백 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 2, 'LITERATURE REVIEW', '국내외 ML 부동산 연구의 세 공백', 5)
    add_text(s, 0.4, 1.25, 12.5, 0.4,
             '서울 부동산 ML 예측 연구는 다수 존재(김이환 2022, 김학현 2023, 조보근 2020, Chun 2025, Kim 2025) — 그러나 다음 세 공백 공유',
             size=12, color=DARK)
    items = [
        ('GAP 1', '공간 단위가 거침',
         'MAUP\n(Modifiable Areal Unit Problem)',
         '자치구(25개) 또는 법정동 단위 집계 변수 사용 → 미시 입지 신호 가려짐'),
        ('GAP 2', '시간 정합성 불완전',
         'Temporal Leakage',
         '환경 변수를 수집 시점 stock으로 분석 기간 전체 병합 → 미래 시설 정보가 과거 거래에 소급 적용'),
        ('GAP 3', '해석이 단일 SHAP에 머묾',
         'Aggregate-only SHAP',
         '권역(강남/비강남)·연도(국면)·면적대(생애주기) 간 설명 신호의 이질성 미포착'),
    ]
    for i, (gap, title, eng, body) in enumerate(items):
        y = 1.85 + i * 1.6
        add_rect(s, 0.4, y, 1.4, 1.4, fill=ACCENT)
        add_text(s, 0.4, y, 1.4, 0.5, gap,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 0.4, y + 0.5, 1.4, 0.9, title,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, 1.95, y, 10.95, 1.4, fill=BG_LIGHT)
        add_text(s, 2.15, y + 0.15, 10.6, 0.45, eng,
                 size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.15, y + 0.7, 10.6, 0.65, body,
                 size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(s, '기존 연구의 세 공백을 정리하면 공간 단위 거침, 시간 정합성 불완전, 단일 SHAP 해석 한계입니다. 본 연구는 이 세 가지를 정면으로 다룹니다.')

    # ===== 06. LITERATURE REVIEW — 본 연구 차별성 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 2, 'LITERATURE REVIEW', '본 연구의 차별성 — 세 공백을 메우는 네 축', 6)
    contribs = [
        ('1. 종속변수 정규화', 'log(㎡당 가격)',
         '규모효과 분리 → 질적·입지 신호 부각'),
        ('2. 공간 정합 개선', '거리 기반 접근성',
         '13 시설군 × 최근접 + 반경 + 도보 분 — MAUP 완화'),
        ('3. 시간 정합 개선', '연도별 활동 스냅샷',
         '학교·학원·어린이집·공원·근린·지하철 6군 시점 정합'),
        ('4. 시공간 이질성', '권역×연도 SHAP',
         '21 서브모델 Top1 변화 추적 — 단일 SHAP 한계 극복'),
    ]
    for i, (t, sub, body) in enumerate(contribs):
        y = 1.5 + i * 1.35
        add_rect(s, 0.4, y, 3.5, 1.15, fill=NAVY)
        add_text(s, 0.5, y + 0.15, 3.3, 0.45, t,
                 size=14, bold=True, color=WHITE)
        add_text(s, 0.5, y + 0.6, 3.3, 0.5, sub,
                 size=14, bold=True, color=WHITE)
        add_rect(s, 4.0, y, 8.9, 1.15, fill=BG_LIGHT)
        add_text(s, 4.2, y + 0.15, 8.6, 0.85, body,
                 size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(s, '본 연구는 종속변수 정규화·공간 정합·시간 정합·시공간 이질성 네 축으로 기존 연구의 세 공백을 메웁니다.')

    # ===== 07. METHODS — 데이터 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 3, 'METHODS', '분석 데이터 — 서울 391,826 거래 (2019.01~2025.12)', 7)
    # 좌측 큰 숫자
    add_rect(s, 0.4, 1.3, 4.5, 5.5, fill=NAVY)
    add_text(s, 0.4, 1.6, 4.5, 0.5, '서울 215개 행정동',
             size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 2.2, 4.5, 1.4, '391,826',
             size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.4, 3.7, 4.5, 0.4, '아파트 매매 실거래 (건)',
             size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 4.5, 4.5, 0.4, '8,601 유니크 단지',
             size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 5.0, 4.5, 0.4, '거래 100% 좌표 커버 (Kakao API)',
             size=12, color=WHITE, align=PP_ALIGN.CENTER)
    # 우측 연도별 표
    add_rect(s, 5.2, 1.3, 7.7, 0.5, fill=NAVY_LT)
    add_text(s, 5.2, 1.3, 7.7, 0.5, '연도별 거래 분포 — 시장 국면을 모두 포함',
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('연도', '건수', '국면', '기준금리 변동'),
        ('2019', '74,896', '안정기', '1.75%'),
        ('2020', '83,916', '유동성 급등 초', '0.50% (저점)'),
        ('2021', '43,379', '급등 말기', '0.50→1.00%'),
        ('2022', '12,788', '금리 충격 조정', '1.00→3.25% (-70% 거래)'),
        ('2023', '35,565', '회복 초기', '3.25→3.50% (고점)'),
        ('2024', '57,710', '회복 확산', '3.50→3.00%'),
        ('2025', '83,572', '안정화', '~2.50%'),
    ]
    for i, row in enumerate(rows):
        y = 1.85 + i * 0.55
        widths = [1.0, 1.5, 2.5, 2.7]
        x = 5.2
        for j, (v, w) in enumerate(zip(row, widths)):
            fill = NAVY_LT if i == 0 else (BG_LIGHT if i % 2 else WHITE)
            add_rect(s, x, y, w, 0.55, fill=fill, line=LIGHT)
            color = WHITE if i == 0 else DARK
            add_text(s, x + 0.05, y, w - 0.1, 0.55, v,
                     size=11, bold=(i == 0 or j == 1), color=color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += w
    add_notes(s, '서울 215개 행정동, 391,826건의 아파트 실거래를 분석했습니다. 카카오 API로 8,601개 단지를 지오코딩해 거래 100% 좌표 커버를 확보했고, 7년간 시장 국면을 모두 포함합니다.')

    # ===== 08. METHODS — 단지 지오코딩 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 3, 'METHODS', '아파트 단지 지오코딩 — Kakao Local API', 8)
    add_text(s, 0.4, 1.4, 12.5, 0.4,
             '국토교통부 실거래 자료에 위경도 정보 부재 → (구·법정동·아파트명) 키로 8,601개 유니크 단지 식별 후 Kakao로 좌표 확보',
             size=12, color=DARK)
    # 단계 박스
    steps = [
        ('STEP 1', '키워드 검색', 'Kakao Local API\n키워드 + 구·법정동',
         '7,999 단지 (93.0%)\n구·법정동 완전 일치(score=3)'),
        ('STEP 2', '쿼리 변형 재시도', '특수문자·로마숫자·접미어 제거\n→ 변형 쿼리 재호출',
         '+602 단지 회복\n→ 누적 100% 커버'),
        ('STEP 3', '좌표체계 통일', '시설별 WGS84 / TM(EPSG:2097) 혼재\n→ pyproj로 EPSG:4326 통일',
         '13 시설군 × 통일 좌표\nBallTree haversine 거리 계산'),
    ]
    for i, (st, t, desc, result) in enumerate(steps):
        x = 0.4 + i * 4.2
        add_rect(s, x, 1.95, 4.0, 0.55, fill=ACCENT)
        add_text(s, x, 1.95, 4.0, 0.55, st,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, 2.5, 4.0, 0.6, fill=NAVY)
        add_text(s, x, 2.5, 4.0, 0.6, t,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, 3.1, 4.0, 1.6, fill=BG_LIGHT)
        add_text(s, x + 0.15, 3.2, 3.7, 1.45, desc,
                 size=11, color=DARK, lines=desc.split('\n'))
        add_rect(s, x, 4.7, 4.0, 1.7, fill=WHITE, line=NAVY)
        add_text(s, x + 0.15, 4.85, 3.7, 1.5, result,
                 size=11, bold=True, color=NAVY, lines=result.split('\n'))
    add_text(s, 0.4, 6.7, 12.5, 0.4,
             '→ 거래 391,826건 100% 매핑 (단지 단위 평균 좌표 또는 법정동 중심좌표 fallback)',
             size=12, bold=True, color=ACCENT)
    add_notes(s, '아파트 단지 지오코딩은 3단계로 진행했습니다. 1차로 93%, 2차 쿼리 변형으로 100% 커버, 3차로 전 시설 좌표체계를 WGS84로 통일했습니다.')

    # ===== 09. METHODS — 거리 기반 변수 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 3, 'METHODS', '거리 기반 입지 변수 — 13 시설군 × 4 축', 9)
    add_text(s, 0.4, 1.4, 12.5, 0.4,
             '단지 좌표 기준 — 최근접 직선거리(m) + 반경 500m/1km/2km 개수 + 도보 추정 시간(분, 직선×1.35/4 km·h)',
             size=12, color=DARK)
    fac_rows = [
        ('교통', '지하철역 (799)', '도보 9.4분 (중앙값)'),
        ('교육', '초·중·고 (1,415) · 학원 (25,437)', '도보 7.3분 (초등)'),
        ('보육', '어린이집 (8,787)', '도보 1.4분 (밀도)'),
        ('근린', '대규모점포 영업중 (37k)', '도보 7.9분 (대형마트)'),
        ('상업', '백화점 (36) · 대형마트 (473)', '도보 38.5분 (백화점)'),
        ('의료', '병원 (1,117) — 종합병원 25 별도', '도보 40.9분 (종합병원)'),
        ('문화', '도서관 (215) · 공원 (130)', '도보 12.0분 (도서관)'),
        ('안전', 'CCTV (129,123) — 500m 밀도', '500m 251개 (밀도)'),
    ]
    add_rect(s, 0.4, 1.95, 12.5, 0.5, fill=NAVY_LT)
    add_text(s, 0.5, 1.95, 4.0, 0.5, '범주', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 4.5, 1.95, 5.0, 0.5, '시설군 (개수)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 9.5, 1.95, 3.4, 0.5, '대표 도보 분(중앙값)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, (cat, fac, walk) in enumerate(fac_rows):
        y = 2.45 + i * 0.5
        fill = WHITE if i % 2 else BG_LIGHT
        add_rect(s, 0.4, y, 12.5, 0.5, fill=fill, line=LIGHT)
        add_text(s, 0.5, y, 4.0, 0.5, cat, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 4.5, y, 5.0, 0.5, fac, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 9.5, y, 3.4, 0.5, walk, size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.4, 6.6, 12.5, 0.5,
             '도서관 12.0분(생활SOC 도보권), 백화점 38.5분(차량권 프리미엄)',
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, '거리 변수는 13 시설군에 대해 최근접 거리, 반경 500m·1km·2km 개수, 도보 추정 분 4축으로 산출했습니다. 도서관 도보권은 12.0분, 백화점은 38.5분으로 차량 권역에 해당합니다.')

    # ===== 10. METHODS — 시점 정합 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 3, 'METHODS', '시점 정합 — 연도별 활동 시설 스냅샷', 10)
    add_text(s, 0.4, 1.4, 12.5, 0.4,
             '거래 연도(Y) 말일(Y-12-31) 기준 활동 중인 시설만 포함 — temporal leakage 제거',
             size=12, color=DARK)
    # 좌측: 시점 정합 6군
    add_rect(s, 0.4, 1.95, 6.2, 0.5, fill=NAVY)
    add_text(s, 0.4, 1.95, 6.2, 0.5, '✓ 시점 정합 적용 시설군 (6)',
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yes_items = [
        '학교 (FOND_YMD 개교일)',
        '학원 (ESTBL_YMD + REG_STTUS_NM)',
        '어린이집 (CRCNFMDT 인가일 + CRABLDT)',
        '공원 (OPEN_YMD 개원일)',
        '근린시설 (APVPERMYMD + DCBYMD)',
        '지하철 (2019~2025 신규 15역 보강)',
    ]
    add_bullets(s, 0.5, 2.55, 6.0, 4.0, yes_items, size=12)
    # 우측: 스냅샷 잔존 5군 (한계)
    add_rect(s, 6.7, 1.95, 6.2, 0.5, fill=ACCENT)
    add_text(s, 6.7, 1.95, 6.2, 0.5, '✗ 스냅샷 잔존 (5) — 한계 명시',
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    no_items = [
        '도서관 (개관일 필드 부재)',
        '백화점 (Kakao 카테고리 — 개점일 부재)',
        '대형마트 (Kakao MT1 — 개점일 부재)',
        'CCTV (설치일 부재 — UPDTDATE만)',
        '병원 (Kakao HP8 — 개원일 부재)',
    ]
    add_bullets(s, 6.8, 2.55, 6.0, 4.0, no_items, size=12)
    add_text(s, 0.4, 6.55, 12.5, 0.5,
             '연도별 시설 활동 수: 학원 14,876(2019) → 25,034(2025), 어린이집 5,818(2019) → 4,106(2025)',
             size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_notes(s, '시점 정합은 6개 시설군에 적용했고, 5개 시설군은 원자료 한계로 2026년 스냅샷을 적용하며 한계로 명시했습니다.')

    # ===== 11. METHODS — 모형·분할 설계 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 3, 'METHODS', '모형·분할 설계 + Ablation 4 시나리오', 11)
    # 위: 모형·분할
    add_rect(s, 0.4, 1.4, 6.2, 0.5, fill=NAVY)
    add_text(s, 0.4, 1.4, 6.2, 0.5, '예측 모형 (3) × 분할 (3)',
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    models = [
        ('OLS', '선형 해석 기준선'),
        ('Random Forest', '배깅 비선형 기준선'),
        ('XGBoost', '그래디언트 부스팅 (대표 해석)'),
    ]
    splits = [
        ('무작위 분할', '같은 단지 중첩 — 가장 쉬움'),
        ('단지 분할', '미관측 단지 외삽 — 가장 엄격'),
        ('시간순 분할', '≤2023 train / ≥2024 test'),
    ]
    for i, (m, d) in enumerate(models):
        y = 2.0 + i * 0.5
        add_rect(s, 0.5, y, 1.7, 0.45, fill=NAVY_LT)
        add_text(s, 0.5, y, 1.7, 0.45, m, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.3, y, 4.2, 0.45, d, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    for i, (sp, d) in enumerate(splits):
        y = 4.0 + i * 0.5
        add_rect(s, 0.5, y, 1.7, 0.45, fill=ACCENT)
        add_text(s, 0.5, y, 1.7, 0.45, sp, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.3, y, 4.2, 0.45, d, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    # 우: Ablation 4 시나리오
    add_rect(s, 6.9, 1.4, 6.0, 0.5, fill=NAVY)
    add_text(s, 6.9, 1.4, 6.0, 0.5, 'Ablation 4 시나리오',
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    abl = [
        ('A', '행정동 집계만', '18 변수 (선행연구 재현)'),
        ('B', '거리 기반 + 시점 무관', '26 변수 (2026 스냅샷)'),
        ('C', '거리 + 시점 정합', '26 변수 (연도별)'),
        ('D', '거리 + 시점 + 행정동', '36 변수 (통합 — 본 연구)'),
    ]
    for i, (k, name, n) in enumerate(abl):
        y = 2.0 + i * 1.05
        add_rect(s, 6.9, y, 0.7, 0.95, fill=NAVY_LT)
        add_text(s, 6.9, y, 0.7, 0.95, k, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 7.7, y + 0.05, 5.2, 0.45, name, size=13, bold=True, color=NAVY)
        add_text(s, 7.7, y + 0.5, 5.2, 0.45, n, size=11, color=DARK)
    add_notes(s, '모형은 OLS·Random Forest·XGBoost 3종을 무작위·단지·시간순 3분할에서 비교했고, Ablation은 4 시나리오로 공간 정합·시간 정합·결합 효과를 분리했습니다.')

    # ===== 12. RESULTS — Ablation (표) =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', 'Ablation — 거리 전환이 시간순 R² +7.1%p 견인', 12)
    # 표
    headers = ['시나리오', '무작위', '시간순', '단지', '변수 수']
    rows = [
        ('A 행정동 집계만', '0.921', '0.800', '0.777', '18'),
        ('B 거리 기반 + 시점 무관', '0.952', '0.856', '0.727', '26'),
        ('C 거리 + 시점 정합', '0.949', '0.851', '0.726', '26'),
        ('D 거리 + 시점 + 행정동', '0.959', '0.871', '0.804', '36'),
    ]
    add_text(s, 0.4, 1.3, 12.5, 0.4, '< 표 > 시나리오별 XGBoost R² (동일 테스트 조건)',
             size=12, bold=True, color=NAVY)
    widths = [4.2, 1.85, 1.85, 1.85, 1.5]
    add_rect(s, 0.4, 1.8, sum(widths), 0.5, fill=NAVY_LT)
    x = 0.4
    for w, h_ in zip(widths, headers):
        add_text(s, x, 1.8, w, 0.5, h_, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w
    for i, row in enumerate(rows):
        y = 2.3 + i * 0.55
        is_best = (row[0].startswith('D '))
        fill = NAVY_LT if is_best else (BG_LIGHT if i % 2 else WHITE)
        x = 0.4
        for j, (w, v) in enumerate(zip(widths, row)):
            add_rect(s, x, y, w, 0.55, fill=fill, line=LIGHT)
            color = WHITE if is_best else (ACCENT if (j > 0 and j < 4 and is_best) else DARK)
            color = WHITE if is_best else DARK
            add_text(s, x + 0.05, y, w - 0.1, 0.55, v,
                     size=12, bold=is_best, color=color,
                     align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += w
    # 핵심 발견 박스
    add_rect(s, 0.4, 5.0, 12.5, 1.6, fill=BG_LIGHT)
    add_text(s, 0.6, 5.1, 12.1, 0.4, '핵심 발견',
             size=14, bold=True, color=ACCENT)
    findings = [
        'A→B (거리 기반 전환) : 시간순 +5.6%p — MAUP 완화 효과의 정량 증거',
        'B→C (시점 정합 단독) : R² 단독 효과 0.5%p 이내 — 성능 아닌 leakage 제거 가치',
        'D (거리+시점+행정동) : 무작위·시간순·단지 분할 모두 최고 — 보완적 결합',
    ]
    add_bullets(s, 0.6, 5.5, 12.1, 1.1, findings, size=12)
    add_notes(s, '''[RESULTS Ablation]
Ablation 결과 거리 전환만으로 시간순 R²가 5.6%p 개선되었고, 시점 정합 단독 효과는 R²로는 작지만 미래 시설 정보 누수를 제거하는 방법론적 가치가 있습니다.
통합 모형 D는 무작위 0.959, 시간순 0.871, 단지 그룹 0.804로 모든 분할에서 최고 성능을 보였습니다.''')

    # ===== 13. RESULTS — Ablation (그림) =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', 'Ablation 시나리오별 XGB R² 시각화', 13)
    add_image(s, os.path.join(FIGS, 'fig10_ablation.png'), 2.4, 1.4, w=8.5)
    add_text(s, 0.4, 6.55, 12.5, 0.5,
             '→ 모든 분할에서 D(통합)가 최고. A→B 거리 전환의 시간순 분할 개선이 가장 큼.',
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, '시각적으로도 D 시나리오가 무작위·시간순·단지 세 분할 모두에서 가장 높은 막대를 보입니다.')

    # ===== 14. RESULTS — 전체 SHAP Top 15 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', '전체 SHAP — 거리 기반 변수가 상위 53% (8/15)', 14)
    add_image(s, os.path.join(FIGS, 'fig4_shap_bar.png'), 0.4, 1.3, w=7.5)
    # 우측 해석 박스
    add_rect(s, 8.3, 1.3, 4.7, 0.5, fill=NAVY_LT)
    add_text(s, 8.3, 1.3, 4.7, 0.5, 'Top 5 변수',
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    top5 = [
        ('1', '강남구분', '12.5%'),
        ('2', '건물연령', '10.4%'),
        ('3', 'M2 통화량', '9.8%'),
        ('4', '전용면적', '7.8%'),
        ('5', '어린이집 1km 내', '6.4%'),
    ]
    for i, (r, name, pct) in enumerate(top5):
        y = 1.85 + i * 0.6
        add_rect(s, 8.3, y, 0.5, 0.55, fill=NAVY)
        add_text(s, 8.3, y, 0.5, 0.55, r, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, 8.85, y, 4.15, 0.55, fill=BG_LIGHT, line=LIGHT)
        add_text(s, 8.95, y, 2.75, 0.55, name, size=12, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 11.5, y, 1.4, 0.55, pct, size=12, bold=True, color=ACCENT,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 8.3, 5.1, 4.7, 1.5,
             '거리 기반 변수가 상위 15 중 8개(53%) — 행정동 단순 집계에서 가려져 있던 미시 입지 신호가 전면에 드러남 (MAUP 완화 효과 실증)',
             size=12, color=DARK)
    add_notes(s, '전체 SHAP Top 5는 강남구분·건물연령·M2·전용면적·어린이집 1km입니다. 거리 기반 변수가 상위 8개 차지합니다.')

    # ===== 15. RESULTS — 비선형 Dependence =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', 'SHAP 비선형 패턴 — 건물연령 U자 / 지하철 임계반응', 15)
    # 차트 width 6.3in × aspect 1.27 → height 4.97in. bottom y = 1.3 + 4.97 = 6.27.
    # 따라서 라벨 텍스트는 6.3 이상으로 내려야 차트와 겹치지 않음.
    add_image(s, os.path.join(FIGS, 'fig6_dep_건물연령.png'), 0.4, 1.3, w=6.3)
    add_image(s, os.path.join(FIGS, 'fig8_dep_subway_nearest.png'), 6.8, 1.3, w=6.3)
    add_text(s, 0.4, 6.35, 6.3, 0.35, '건물연령 25~30년 반등 — 재건축 기대 신호',
             size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 6.8, 6.35, 6.3, 0.35, '지하철 500m 이내 강한 양 → 평탄화 — 도보 역세권 프리미엄',
             size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 6.8, 12.5, 0.4,
             '→ OLS 선형 가정으로 포착되지 않는 비선형 패턴을 SHAP이 정량 시각화',
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, '건물연령은 25~30년 구간에서 반등하는 U자형이고, 지하철 최근접거리는 500m 이내 강한 양에서 평탄화되는 임계반응형입니다. OLS는 못 잡는 패턴입니다.')

    # ===== 16. RESULTS — 권역별 SHAP =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', '권역별 SHAP — 강남(상업·안전) vs 비강남(사교육·주거)', 16)
    add_image(s, os.path.join(FIGS, 'fig11_region_shap.png'), 0.4, 1.3, w=8.0)
    # 우측 표
    add_rect(s, 8.7, 1.3, 4.4, 0.5, fill=NAVY_LT)
    add_text(s, 8.7, 1.3, 4.4, 0.5, '권역별 Top 5 변수',
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 8.7, 1.85, 2.2, 0.4, '강남3구', size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 10.9, 1.85, 2.2, 0.4, '비강남', size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    g = ['건물연령', '백화점 최근접', '전용면적', '백화점수(행정동)', 'CCTV 500m']
    n = ['건물연령', '학원수(행정동)', '전용면적', '어린이집수(행정동)', '층']
    for i, (gv, nv) in enumerate(zip(g, n)):
        y = 2.3 + i * 0.5
        add_rect(s, 8.7, y, 2.2, 0.45, fill=BG_LIGHT, line=LIGHT)
        add_text(s, 8.7, y, 2.2, 0.45, f'{i+1}. {gv}', size=11, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, 10.9, y, 2.2, 0.45, fill=BG_LIGHT, line=LIGHT)
        add_text(s, 10.9, y, 2.2, 0.45, f'{i+1}. {nv}', size=11, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.4, 6.5, 12.5, 0.5,
             '→ R² 강남 0.964, 비강남 0.952 — 두 권역의 가격 설명 구조에 질적 차이',
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, '강남은 백화점·CCTV 등 상업·안전 신호, 비강남은 학원·어린이집 등 사교육·보육 신호가 상위에 있습니다.')

    # ===== 17. RESULTS — 연도×권역 R² 히트맵 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', '연도×권역 R² — 2022 조정기 동시 dip', 17)
    add_image(s, os.path.join(FIGS, 'fig12_year_region_heatmap.png'), 1.5, 1.3, w=6.0)
    # 우측 해석
    add_rect(s, 8.0, 1.3, 5.0, 0.5, fill=NAVY_LT)
    add_text(s, 8.0, 1.3, 5.0, 0.5, '핵심 패턴',
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    insights = [
        '두 권역 모두 0.92 이상 안정\n(2021 비강남 0.919만 예외)',
        '2022 R² 동시 0.857 dip\n→ 거래 −70% + 금리 1.25→3.25% 급등',
        '전체 R² 0.892 > 권역별 0.857\n→ 권역 간 가격 수준 차로 pooled variance ↑',
    ]
    for i, ins in enumerate(insights):
        y = 1.95 + i * 1.55
        add_rect(s, 8.0, y, 5.0, 1.4, fill=BG_LIGHT)
        add_text(s, 8.2, y + 0.1, 4.6, 1.25, ins,
                 size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE,
                 lines=ins.split('\n'))
    add_notes(s, '2022년에 두 권역 모두 R² 0.857로 dip을 보입니다. 거래량 -70%와 금리 급등이 동반된 조정기에서 예측 안정성이 약화된 것으로 해석합니다.')

    # ===== 18. RESULTS — 연도별 Top1 변화 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 4, 'RESULTS', '연도별 Top 1 SHAP 변수 — 강남 6회 교체 vs 비강남 안정', 18)
    add_image(s, os.path.join(FIGS, 'fig13_top1_timeline.png'), 0.92, 1.3, w=11.5)
    # 하단 핵심 메시지
    add_rect(s, 0.4, 5.5, 12.5, 1.7, fill=BG_LIGHT)
    add_text(s, 0.6, 5.6, 12.1, 0.4, '핵심 발견 — 권역별 시간 안정성의 차이',
             size=14, bold=True, color=ACCENT)
    add_bullets(s, 0.6, 6.0, 12.1, 1.2, [
        '강남3구 : 백화점수→건물연령→CCTV→건물연령→종합병원→건물연령→백화점 최근접 — 7년간 6회 교체',
        '비강남 : 7년 중 5년 건물연령 1위 (2020 어린이집수·2021 전용면적 예외) — 안정적 구조',
        '단일 SHAP에서 보이지 않던 "권역별 시간 안정성 차이"가 본 연구의 차별 발견',
    ], size=11)
    add_notes(s, '강남 Top1은 7년간 6회 바뀌고, 비강남은 거의 매년 건물연령이 1위입니다. 두 권역의 가격 설명 구조가 서로 다른 시간 안정성을 가집니다.')

    # ===== 19. CONCLUSION — 기여 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 5, 'CONCLUSION', '학술·실무 기여', 19)
    # 한 줄 기여
    add_rect(s, 0.4, 1.3, 12.5, 0.95, fill=NAVY)
    add_text(s, 0.6, 1.4, 12.1, 0.45, '본 연구의 한 줄 기여',
             size=13, color=WHITE)
    add_text(s, 0.6, 1.7, 12.1, 0.5,
             '"공간 단위·시간 정합성을 통제한 뒤 예측 신호의 권역×연도 이질성을 설명 가능한 방식으로 검증"',
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 학술 기여
    add_rect(s, 0.4, 2.55, 6.2, 0.5, fill=NAVY_LT)
    add_text(s, 0.4, 2.55, 6.2, 0.5, '학술적 기여', size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, 0.5, 3.15, 6.0, 3.5, [
        'MAUP 완화의 실증 — 시간순 R² +7.1%p 개선',
        '시간 정합성의 방법론적 통제 — leakage 제거',
        '시공간 이질성의 XAI 해석 — 권역×연도 21 SHAP',
        '발견(insight)형 기여 (algorithm·module 아님)',
    ], size=12)
    # 실무 기여
    add_rect(s, 6.7, 2.55, 6.2, 0.5, fill=ACCENT)
    add_text(s, 6.7, 2.55, 6.2, 0.5, '실무·정책 참고점', size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, 6.8, 3.15, 6.0, 3.5, [
        'AVM 설명가능성 강화 (SHAP 분해)',
        '권역 맞춤 모형 — 강남/비강남 분리 운용',
        '국면 전환기 진단 — 2022 dip 사례',
        '재건축 기대의 비선형 신호 (건물연령 U자)',
    ], size=12)
    add_notes(s, '학술 기여 4가지와 실무 기여 4가지를 정리했습니다.')

    # ===== 20. CONCLUSION — 한계 + 마무리 =====
    s = prs.slides.add_slide(blank)
    add_chapter_header(s, 5, 'CONCLUSION', '연구의 한계 · 후속 연구 · 마무리', 20)
    # 한계
    add_rect(s, 0.4, 1.3, 6.2, 0.5, fill=ACCENT)
    add_text(s, 0.4, 1.3, 6.2, 0.5, '한계',
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, 0.5, 1.9, 6.0, 3.0, [
        '5 시설군 스냅샷 잔존 (도서관·백화점·대형마트·CCTV·병원)',
        '도보 추정은 직선거리 × 1.35 / 4 km·h 근사',
        '거시·정책 변수 미포함 (LTV·DTI·재건축 인허가)',
        'SHAP은 인과 효과 아닌 예측 기여도',
        '단지 분할 외삽 한계 + 비교 모형 범위 제한',
    ], size=11)
    # 후속 연구
    add_rect(s, 6.7, 1.3, 6.2, 0.5, fill=NAVY_LT)
    add_text(s, 6.7, 1.3, 6.2, 0.5, '후속 연구',
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, 6.8, 1.9, 6.0, 3.0, [
        '정밀 보행 네트워크 거리 (OSRM·ORS) 검증',
        '시설 개폐업 완전 패널 + 정책 이벤트 결합',
        '반복매매 지수와의 교차검증',
        'LightGBM·CatBoost·TabNet과의 비교',
        '수도권·지방 대도시 외적 타당성 검증',
    ], size=11)
    # 감사
    add_text(s, 0.4, 6.4, 12.5, 0.5, '감사합니다 · Q & A',
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_notes(s, '한계와 후속 연구 방향을 정리하며 마무리합니다. 감사합니다. 질문 받겠습니다.')

    prs.save(OUT)
    sz = os.path.getsize(OUT) // 1024
    print(f'✅ {OUT} ({sz}KB, {len(prs.slides)} 슬라이드, HUSA 학술 표준 디자인)')


if __name__ == '__main__':
    main()
