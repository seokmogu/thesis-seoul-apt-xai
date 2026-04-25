#!/usr/bin/env python3
"""
중간발표 5분 PPT (16:9, 한글 폰트 cascade).

8 슬라이드: 표지 / Why / 방법 / 어블레이션 / 권역×연도 SHAP /
            연도별 변화 / 한계·후속 / 마무리.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, 'figures')
OUT = os.path.join(ROOT, 'paper', '중간발표_박현근.pptx')

# 한글 폰트 cascade — Apple SD Gothic Neo(Mac) / Malgun Gothic(Win) / Nanum Gothic(Linux)
KR = 'Apple SD Gothic Neo'
KR_FALLBACK = 'Malgun Gothic'

PRIMARY = RGBColor(0x1F, 0x77, 0xB4)   # 파랑 (학술)
ACCENT = RGBColor(0xD6, 0x27, 0x28)    # 강남 빨강
DARK = RGBColor(0x39, 0x39, 0x39)
LIGHT = RGBColor(0x7A, 0x7A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF4, 0xF4, 0xF4)


def set_font(run, size=18, bold=False, color=DARK, font=KR):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한글 동아시아 폰트 명시
    rPr = run._r.get_or_add_rPr()
    for ex in rPr.findall(qn('a:ea')):
        rPr.remove(ex)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    # 라틴 폰트도 동일 (한국어 ppt에서 일관)
    for ex in rPr.findall(qn('a:latin')):
        rPr.remove(ex)
    la = etree.SubElement(rPr, qn('a:latin'))
    la.set('typeface', font)


def add_text(slide, left_in, top_in, w_in, h_in, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, lines=None):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if lines is None:
        lines = [text]
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = l
        set_font(run, size, bold, color)
    return tb


def add_bullets(slide, left_in, top_in, w_in, h_in, bullets, size=18, color=DARK):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = '• ' + b
        set_font(run, size, False, color)
    return tb


def add_rect(slide, left_in, top_in, w_in, h_in, fill_color=BG_LIGHT):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in),
                                 Inches(w_in), Inches(h_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_image(slide, path, left_in, top_in, w_in=None, h_in=None):
    if not os.path.exists(path):
        return None
    kw = {}
    if w_in:
        kw['width'] = Inches(w_in)
    if h_in:
        kw['height'] = Inches(h_in)
    return slide.shapes.add_picture(path, Inches(left_in), Inches(top_in), **kw)


def add_notes(slide, script):
    """슬라이드 발표자 노트(speaker notes) 추가."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = script
    # 노트 폰트도 한글 명시
    for p in tf.paragraphs:
        for r in p.runs:
            try:
                set_font(r, size=12, bold=False, color=DARK)
            except Exception:
                pass


def add_header(slide, num, title):
    """슬라이드 상단 헤더 — 학술 보수 스타일 (검정 텍스트, 하단 가는 회색 선)."""
    add_text(slide, 0.4, 0.15, 1.2, 0.55, f'{num:02d}', size=20, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.2, 0.15, 11.5, 0.55, title, size=22, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    # 하단 구분선 (가는 회색)
    add_rect(slide, 0.4, 0.78, 12.5, 0.02, fill_color=LIGHT)


def add_footer(slide):
    add_text(slide, 0.4, 7.05, 12.5, 0.4, '한양대학교 부동산융합대학원 · 박현근 · 2026',
             size=11, color=LIGHT, align=PP_ALIGN.LEFT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # =========== 1. 표지 ===========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill_color=WHITE)
    add_rect(s, 0, 6.5, 13.333, 1.0, fill_color=BG_LIGHT)
    add_text(s, 0.7, 1.5, 12.0, 0.6, '한양대학교 부동산융합대학원 석사학위논문 (중간발표)',
             size=18, color=LIGHT)
    add_text(s, 0.7, 2.3, 12.0, 1.6,
             'XGBoost와 SHAP을 활용한\n서울시 아파트 단위면적당 매매가격의 설명 패턴 분석',
             size=34, bold=True, color=DARK,
             lines=['XGBoost와 SHAP을 활용한',
                    '서울시 아파트 단위면적당 매매가격의 설명 패턴 분석'])
    add_text(s, 0.7, 4.5, 12.0, 0.6,
             '— 거리 기반 입지 변수 · 연도별 시점 정합 · 권역×연도 SHAP 분해 —',
             size=18, color=DARK)
    add_text(s, 0.7, 5.5, 12.0, 0.5, '도시부동산정책전공  박현근',
             size=20, color=DARK)
    add_text(s, 0.7, 6.05, 12.0, 0.4, '지도교수  고준호', size=18, color=DARK)
    add_text(s, 0.7, 6.65, 12.0, 0.7, '2026. 04', size=14, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(s, '''[표지 · 약 30초]
안녕하세요. 한양대학교 부동산융합대학원 석사 과정 박현근입니다.
"XGBoost와 SHAP을 활용한 서울시 아파트 단위면적당 매매가격의 설명 패턴 분석"이라는 주제로 중간발표를 진행하겠습니다.
본 연구는 거리 기반 입지 변수, 연도별 시점 정합, 그리고 권역×연도 SHAP 분해 세 축을 결합한 분석입니다. 5분 정도 핵심만 말씀드리겠습니다.''')

    # =========== 2. 연구 배경 (Why) ===========
    s = prs.slides.add_slide(blank)
    add_header(s, 1, '연구 배경 — 왜 이 주제인가')
    add_text(s, 0.7, 1.0, 12.0, 0.5, '세 가지 층위의 동기',
             size=22, bold=True, color=DARK)
    box_y = 1.7
    for i, (title, body) in enumerate([
        ('실무·제도', 'AVM 확산 가속, 감정평가·은행·세무에 “왜 이 가격인가”를\n설명할 해석 프레임워크 요구'),
        ('학술 공백', '국내 ML 부동산 연구의 세 공백:\n공간 단위 거침(MAUP) · 시간역전 누수 · 단일 SHAP'),
        ('시장 국면', '2019~2025: 유동성 급등(2020-21) → 금리 충격 조정(2022)\n→ 완만한 회복(2023-25)을 단일 기간에 포함'),
    ]):
        x = 0.7 + i * 4.1
        add_rect(s, x, box_y, 3.9, 4.0, fill_color=BG_LIGHT)
        add_text(s, x + 0.2, box_y + 0.2, 3.5, 0.6, title,
                 size=20, bold=True, color=DARK)
        lines = body.split('\n')
        add_text(s, x + 0.2, box_y + 1.0, 3.5, 2.8, body,
                 size=15, color=DARK, lines=lines)
    add_text(s, 0.7, 6.2, 12.0, 0.6,
             '→ Q1 공간 정합 · Q2 시간 정합 · Q3 시공간 이질성 · Q4 해석 프레임워크',
             size=16, bold=True, color=ACCENT)
    add_footer(s)
    add_notes(s, '''[연구 배경 · 약 50초]
본 연구의 동기는 세 층위입니다.
첫째, 실무·제도 층위입니다. 자동감정평가 AVM의 실무 확산이 가속화되면서, "왜 이 가격인가"를 재현 가능하게 설명할 해석 프레임워크가 요구되고 있습니다.
둘째, 학술 층위입니다. 기존 ML 부동산 연구는 자치구 단위의 거친 공간 단위, 시간역전 정보누수, 단일 SHAP 해석에 머무르는 세 공백을 공유합니다.
셋째, 시장 국면 층위입니다. 2019년부터 2025년까지의 7년은 유동성 주도 급등, 금리 충격 조정, 완만한 회복을 단일 기간 안에 포함합니다.
이로부터 공간 정합·시간 정합·시공간 이질성·해석 프레임워크라는 네 가지 연구 질문을 도출했습니다.''')

    # =========== 3. 데이터·방법 ===========
    s = prs.slides.add_slide(blank)
    add_header(s, 2, '데이터와 방법 — 거리 기반 + 시점 정합')
    add_text(s, 0.7, 1.0, 12.0, 0.5, '서울 215개 행정동 · 391,826 거래 (2019.01~2025.12)',
             size=20, bold=True, color=DARK)
    add_bullets(s, 0.7, 1.7, 12.0, 5.0, [
        '아파트 단지 8,601개 Kakao Local API 지오코딩 (거래건수 기준 100% 커버)',
        '13 시설군 거리 기반 변수: 최근접거리(m) + 반경 500m/1km/2km 개수 + 도보 추정 시간(분)',
        '시점 정합: 학교·학원·어린이집·공원·근린시설·지하철 6군에 대해 연도별 활동 스냅샷',
        '종속변수: log(㎡당 거래가격) — 규모효과 정규화 (교수 피드백 반영)',
        '모형: OLS·Random Forest·XGBoost — 무작위·단지 Group·시간순 세 분할 검증',
        'SHAP 분해: 전체·권역(강남3구/비강남)·연도·연도×권역(21 서브모델)',
        '어블레이션 4 시나리오: A 행정동만 / B 거리·시점무관 / C 거리+시점 / D 통합',
    ], size=15)
    add_footer(s)
    add_notes(s, '''[데이터와 방법 · 약 60초]
분석 대상은 서울 215개 행정동의 39만 1천여 거래입니다.
8천 6백여 개 아파트 단지를 카카오 Local API로 지오코딩해 거래 100% 좌표 커버리지를 확보했습니다.
13개 시설군에 대해 최근접 직선거리, 반경 500m·1km·2km 개수, 그리고 도보 추정 시간 네 축을 산출했습니다.
시점 정합은 학교, 학원, 어린이집, 공원, 근린시설, 지하철 6개 군에 대해 거래 연도별 활동 시설만 필터링한 스냅샷을 만들었습니다.
종속변수는 단위면적당 가격에 자연로그를 취해 규모효과를 정규화했고, 이는 교수님 피드백을 직접 반영한 부분입니다.
모형은 OLS, 랜덤포레스트, XGBoost를 무작위·단지 그룹·시간순 세 분할에서 비교했고, SHAP은 전체·권역·연도·연도×권역의 네 단면으로 분해했습니다.
어블레이션은 행정동만, 거리만 시점무관, 거리+시점, 거리+시점+행정동 통합 네 시나리오로 설계했습니다.''')

    # =========== 4. 어블레이션 결과 ===========
    s = prs.slides.add_slide(blank)
    add_header(s, 3, '어블레이션 — 거리 전환이 시간순 R² +7.1%p 개선 견인')
    add_image(s, os.path.join(FIGS, 'fig10_ablation.png'), 0.5, 1.0, w_in=7.8)
    add_text(s, 8.6, 1.2, 4.4, 0.5, '핵심 발견', size=20, bold=True, color=DARK)
    add_bullets(s, 8.6, 1.8, 4.4, 5.0, [
        'A→B 거리 변환만으로 시간순 +5.6%p',
        'B→C 시점 정합 단독 R² 효과는 작음 (~0.5%p)\n  → 성능이 아니라 미래 시점 정보의\n     누수 제거(방법론적 타당성)',
        'D 통합(거리+시점+행정동)이 무작위·시간순·\nGroup 세 분할 모두에서 최고 성능',
        '단지 Group: A 0.777 → B 0.727(일시 하락)\n→ D 0.804(거리·행정동 보완)',
    ], size=14)
    add_footer(s)
    add_notes(s, '''[어블레이션 결과 · 약 50초]
어블레이션 분석에서 시나리오 A는 행정동 집계만, B는 거리 변수에 시점 무관, C는 거리에 시점 정합, D는 셋을 모두 통합한 모형입니다.
첫 번째 핵심 결과는, 거리 변환 단독으로 시간순 R²가 5.6 퍼센트포인트 개선됐다는 것입니다. 즉 MAUP 완화 효과가 가장 큰 동인입니다.
두 번째는, 시점 정합 단독 R² 효과는 0.5 퍼센트포인트 이내로 매우 작습니다. 이는 시점 정합이 R² 개선이 아니라 미래 시설 정보 누수를 제거하는 방법론적 타당성 확보에 기여한다는 의미로 해석합니다.
세 번째는 통합 D 모형이 무작위, 시간순, 단지 그룹 모든 분할에서 최고 성능을 보였다는 점입니다. 시간순 분할 R²는 0.871로 행정동 집계 대비 7.1 퍼센트포인트 개선됐습니다.''')

    # =========== 5. 권역·연도×권역 SHAP ===========
    s = prs.slides.add_slide(blank)
    add_header(s, 4, '권역×연도 SHAP — 강남 변동 vs 비강남 안정')
    add_image(s, os.path.join(FIGS, 'fig12_year_region_heatmap.png'), 0.5, 1.0, w_in=5.5)
    add_image(s, os.path.join(FIGS, 'fig13_top1_timeline.png'), 6.5, 1.0, w_in=6.5)
    add_text(s, 0.5, 4.6, 6.0, 0.5, 'R² 7년×3권역', size=16, bold=True, color=DARK)
    add_bullets(s, 0.5, 5.0, 6.0, 2.0, [
        '두 권역 모두 0.92 이상 안정 (2022 dip 0.857)',
        '2022 거래 70% 감소 + 금리 1.25→3.25% 급등',
    ], size=13)
    add_text(s, 6.5, 4.6, 6.5, 0.5, 'Top 1 연도별 변동', size=16, bold=True, color=DARK)
    add_bullets(s, 6.5, 5.0, 6.5, 2.0, [
        '강남: 백화점수→건물연령→CCTV→건물연령\n→종합병원→건물연령→백화점 최근접 (6회 교체)',
        '비강남: 건물연령 일관(2020 어린이집 1회)\n— 7년 안정',
    ], size=13)
    add_footer(s)
    add_notes(s, '''[권역×연도 SHAP · 약 60초]
왼쪽은 연도×권역 R² 히트맵입니다. 두 권역 모두 2022년에 0.857로 동시에 dip을 보이고, 다른 해는 대체로 0.92 이상으로 안정적입니다.
2022년의 동시 저하는 거래량 70% 감소와 기준금리 1.25에서 3.25 퍼센트 급등이 동반된 조정기 영향으로 해석됩니다.
오른쪽은 연도별 Top 1 변수 변화입니다.
강남3구의 Top 1 변수는 백화점수, 건물연령, CCTV, 건물연령, 종합병원 최근접, 건물연령, 백화점 최근접 순으로 7년간 6번 바뀝니다.
반면 비강남은 거의 매년 건물연령이 1위를 차지해 7년 안정적인 구조를 보입니다.
이는 두 권역의 가격 설명 신호가 서로 다른 시간 안정성을 가진다는 의미이며, 단일 SHAP 분석에서는 보이지 않던 새로운 발견입니다.''')

    # =========== 6. SHAP 글로벌 + 비선형 ===========
    s = prs.slides.add_slide(blank)
    add_header(s, 5, '전체 SHAP과 비선형 패턴 — 거리 변수가 상위 50%')
    add_image(s, os.path.join(FIGS, 'fig4_shap_bar.png'), 0.5, 1.0, w_in=6.2)
    add_image(s, os.path.join(FIGS, 'fig6_dep_건물연령.png'), 7.0, 1.0, w_in=6.0)
    add_text(s, 0.5, 5.7, 6.2, 0.5, 'Top 5: 강남 12.5 · 건물연령 10.4 · M2 9.8 · 면적 7.8 · 어린이집 1km 6.4',
             size=12, color=DARK)
    add_text(s, 7.0, 5.7, 6.0, 0.5, '건물연령 25~30년 구간 U자형 — 재건축 기대 신호',
             size=12, color=DARK)
    add_text(s, 0.7, 6.4, 12.0, 0.5,
             '→ 상위 15 중 8개(53%)가 거리 기반 변수: MAUP 완화 효과 실증',
             size=15, bold=True, color=ACCENT)
    add_footer(s)
    add_notes(s, '''[전체 SHAP과 비선형 · 약 30초]
전체 모형 SHAP Top 5는 강남구분 12.5%, 건물연령 10.4%, M2 통화량 9.8%, 전용면적 7.8%, 어린이집 1km 내 개수 6.4%입니다.
상위 15개 중 8개, 53%가 거리 기반 변수입니다. 행정동 단순 집계 방식에서 가려져 있던 미시 입지 신호가 거리 기반 설계에서 전면에 드러난 것이며, MAUP 완화 효과의 실증입니다.
오른쪽은 건물연령 SHAP Dependence Plot입니다. 25에서 30년 구간에서 반등하는 U자형 비선형 패턴은 재건축 기대 신호로, OLS 선형 가정으로는 포착되지 않는 부분입니다.''')

    # =========== 7. 한계와 후속 ===========
    s = prs.slides.add_slide(blank)
    add_header(s, 6, '한계와 후속 연구 — 정직한 진단')
    add_text(s, 0.7, 1.0, 6.0, 0.5, '한계', size=20, bold=True, color=ACCENT)
    add_bullets(s, 0.7, 1.6, 6.0, 5.0, [
        '도서관·백화점·대형마트·CCTV·병원 5군은 2026 스냅샷',
        '도보 시간은 직선거리 × 1.35 / 4km/h 근사 (네트워크 라우팅 아님)',
        '거시·정책 변수(LTV·DTI·DSR·재건축 인허가) 미포함',
        'SHAP은 인과 효과가 아닌 예측 기여도',
        'Group 분할 R² 하락 — 미관측 신규 단지 외삽 한계',
    ], size=14)
    add_text(s, 7.0, 1.0, 6.0, 0.5, '후속 연구', size=20, bold=True, color=DARK)
    add_bullets(s, 7.0, 1.6, 6.0, 5.0, [
        '정밀 보행 네트워크 거리 (OSRM/ORS/상용 API) 민감도 분석',
        '시설 개폐업 완전 패널 + 정책 이벤트 결합',
        '반복매매 지수와의 교차검증',
        '딥러닝(TabNet, Transformer)·LIME/ALE/PI 비교',
        '수도권·지방 대도시로 확장 (외적 타당성)',
    ], size=14)
    add_footer(s)
    add_notes(s, '''[한계와 후속 연구 · 약 30초]
한계는 다섯 가지입니다.
첫째, 도서관·백화점·대형마트·CCTV·병원 5개 군이 2026년 스냅샷이라 시점 정합이 미완성입니다.
둘째, 도보 시간은 직선거리 1.35배 보정 근사이며 실제 보행 네트워크 거리가 아닙니다.
셋째, 거시·정책 변수(LTV·DTI·DSR·재건축 인허가)가 미포함입니다.
넷째, SHAP은 인과 효과가 아닌 예측 기여도입니다.
다섯째, 단지 그룹 분할 R²가 낮아 미관측 신규 단지 외삽 한계가 있습니다.
후속 연구로는 정밀 보행 네트워크 거리 검증, 시설 개폐업 완전 패널, 정책 이벤트 결합, 딥러닝 비교, 외적 타당성 검증을 제시합니다.''')

    # =========== 8. 마무리 ===========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill_color=WHITE)
    add_rect(s, 0, 0.78, 12.5, 0.02, fill_color=LIGHT)  # 가는 회색 선만
    add_text(s, 0.7, 1.5, 12.0, 1.0, '본 연구의 기여 한 줄 요약',
             size=22, bold=True, color=LIGHT)
    add_text(s, 0.7, 2.5, 12.0, 1.6,
             '“공간 단위·시간 정합성을 통제한 뒤,\n예측 신호의 권역×연도 이질성을\n설명 가능한 방식으로 검증한 연구”',
             size=30, bold=True, color=DARK,
             lines=['"공간 단위·시간 정합성을 통제한 뒤,',
                    '예측 신호의 권역×연도 이질성을',
                    '설명 가능한 방식으로 검증한 연구"'])
    add_text(s, 0.7, 5.0, 12.0, 0.5, '교수 피드백 정량 답변',
             size=18, bold=True, color=DARK)
    add_bullets(s, 0.7, 5.6, 12.0, 1.4, [
        '도서관 도보 → 중앙값 12.0분 (생활SOC 도보권)',
        '백화점 조건 → 중앙값 38.5분 (도보권 밖, 차량권 프리미엄)',
        '시간 변동 → 6 시설군 연도별 활동 스냅샷, 5 시설군은 한계로 명시',
    ], size=14)
    add_text(s, 0.7, 6.9, 12.0, 0.4, '감사합니다 · Q&A',
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_notes(s, '''[마무리 · 약 30초]
본 연구의 기여를 한 줄로 정리하면, 공간 단위와 시간 정합성을 통제한 뒤 예측 신호의 권역×연도 이질성을 설명 가능한 방식으로 검증한 연구입니다.
교수님 피드백에 대한 정량 답변으로 도서관 도보 중앙값은 12.0분으로 생활SOC 도보권에 해당하고, 백화점은 38.5분으로 도보권 밖, 시간 변동은 6개 시설군에 연도별 활동 스냅샷을 적용해 다뤘습니다.
이상으로 발표를 마치겠습니다. 감사합니다. 질문 받겠습니다.''')

    prs.save(OUT)
    sz = os.path.getsize(OUT) // 1024
    print(f'✅ {OUT} ({sz}KB, {len(prs.slides)} 슬라이드)')


if __name__ == '__main__':
    main()
